from __future__ import annotations

import json
import mimetypes
import os
import re
import threading
import time
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable


# The registry only contains files explicitly selected through AgentStudio's
# native file picker.  API callers receive opaque ids instead of permission to
# read arbitrary local paths.
_REGISTRY_TTL_SECONDS = 12 * 60 * 60
_MAX_REGISTRY_ITEMS = 200
_MAX_FILE_BYTES = 20 * 1024 * 1024
_MAX_FILE_CONTEXT_CHARS = 36_000
_DEFAULT_TOTAL_CONTEXT_CHARS = 90_000


_SECRET_ASSIGNMENT_RE = re.compile(
    r"(?im)^(?P<prefix>\s*(?:export\s+)?(?:[A-Z0-9_]*(?:API_KEY|TOKEN|SECRET|PASSWORD|PASSWD|PWD)|DATABASE_URL|DB_URL|SUPABASE_DB_PASSWORD|PGPASSWORD)\s*=\s*)(?P<value>[^\r\n]+)$"
)
_SECRET_TOKEN_PATTERNS = (
    re.compile(r"\bsk-[A-Za-z0-9_\-]{16,}\b"),
    re.compile(r"\bAIza[0-9A-Za-z_\-]{20,}\b"),
    re.compile(r"\bgh[pousr]_[A-Za-z0-9]{20,}\b"),
    re.compile(r"\bxox[baprs]-[A-Za-z0-9-]{16,}\b"),
)
_DATABASE_PASSWORD_RE = re.compile(
    r"(?i)(?P<prefix>(?:postgres(?:ql)?|mysql|mariadb|redis)://[^:\s/@]+:)(?P<password>[^@\s/]+)(?P<suffix>@)"
)


def redact_sensitive_text(value: str) -> str:
    """Mask credentials before selected-file contents enter any AI context.

    Reference files are user-selected analysis inputs, not a secret transport.
    Keep variable names/connection structure so the model can understand the
    architecture while preventing actual API keys, passwords and tokens from
    being echoed into prompts, chat history or generated workflow artifacts.
    """
    text = str(value or '')
    if not text:
        return text

    def replace_assignment(match: re.Match) -> str:
        raw = str(match.group('value') or '').strip()
        lowered = raw.casefold()
        placeholders = (
            'os.getenv', 'getenv(', '${', '$env:', '%', '<', 'your_',
            'changeme', 'example', 'placeholder', '[redacted]', '***',
        )
        if any(token in lowered for token in placeholders):
            return match.group(0)
        return f"{match.group('prefix')}[REDACTED]"

    text = _SECRET_ASSIGNMENT_RE.sub(replace_assignment, text)
    text = _DATABASE_PASSWORD_RE.sub(
        lambda m: f"{m.group('prefix')}[REDACTED]{m.group('suffix')}",
        text,
    )
    for pattern in _SECRET_TOKEN_PATTERNS:
        text = pattern.sub('[REDACTED_TOKEN]', text)
    return text


_TEXT_EXTENSIONS = {
    '.txt', '.md', '.markdown', '.rst', '.log', '.csv', '.tsv',
    '.json', '.jsonl', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf',
    '.xml', '.html', '.htm', '.css', '.scss', '.sass', '.less',
    '.py', '.pyi', '.js', '.jsx', '.ts', '.tsx', '.mjs', '.cjs',
    '.java', '.kt', '.kts', '.cs', '.c', '.h', '.cpp', '.hpp', '.cc',
    '.go', '.rs', '.rb', '.php', '.swift', '.dart', '.scala',
    '.sh', '.bash', '.zsh', '.ps1', '.psm1', '.bat', '.cmd',
    '.sql', '.graphql', '.gql', '.proto', '.dockerfile', '.env',
    '.gitignore', '.editorconfig', '.properties', '.gradle', '.sln', '.csproj',
    '.vb', '.vbs', '.vue', '.svelte', '.lua', '.r', '.m', '.tex',
}


@dataclass
class AttachmentRecord:
    attachment_id: str
    path: str
    name: str
    size: int
    extension: str
    mime_type: str
    registered_at: float
    project_relative_path: str = ''


_lock = threading.RLock()
_registry: dict[str, AttachmentRecord] = {}
_content_cache: dict[str, dict] = {}


def _cleanup_locked() -> None:
    now = time.time()
    expired = [
        key for key, value in _registry.items()
        if now - value.registered_at > _REGISTRY_TTL_SECONDS
    ]
    for key in expired:
        _registry.pop(key, None)
        _content_cache.pop(key, None)

    if len(_registry) <= _MAX_REGISTRY_ITEMS:
        return
    oldest = sorted(_registry.values(), key=lambda item: item.registered_at)
    for item in oldest[: max(0, len(_registry) - _MAX_REGISTRY_ITEMS)]:
        _registry.pop(item.attachment_id, None)
        _content_cache.pop(item.attachment_id, None)


def _project_relative(path: Path, project_root: str) -> str:
    if not project_root:
        return ''
    try:
        root = Path(project_root).expanduser().resolve()
        return path.relative_to(root).as_posix()
    except Exception:
        return ''


def _record_to_dict(record: AttachmentRecord) -> dict:
    return {
        'attachment_id': record.attachment_id,
        'path': record.path,
        'name': record.name,
        'size': record.size,
        'extension': record.extension,
        'mime_type': record.mime_type,
        'project_relative_path': record.project_relative_path,
        'registered_at': record.registered_at,
    }


def register_selected_files(paths: Iterable[str], project_root: str = '') -> list[dict]:
    registered: list[dict] = []
    seen_paths: set[str] = set()

    with _lock:
        _cleanup_locked()
        for raw in paths:
            if not str(raw or '').strip():
                continue
            path = Path(str(raw)).expanduser().resolve()
            normalized = os.path.normcase(str(path))
            if normalized in seen_paths or not path.exists() or not path.is_file():
                continue
            seen_paths.add(normalized)

            try:
                size = int(path.stat().st_size)
            except OSError:
                continue
            if size > _MAX_FILE_BYTES:
                registered.append({
                    'ok': False,
                    'path': str(path),
                    'name': path.name,
                    'size': size,
                    'message': f'파일이 너무 큽니다. 최대 {_MAX_FILE_BYTES // (1024 * 1024)}MB까지 등록할 수 있습니다.',
                })
                continue

            attachment_id = uuid.uuid4().hex
            mime_type = mimetypes.guess_type(path.name)[0] or 'application/octet-stream'
            record = AttachmentRecord(
                attachment_id=attachment_id,
                path=str(path),
                name=path.name,
                size=size,
                extension=path.suffix.casefold(),
                mime_type=mime_type,
                registered_at=time.time(),
                project_relative_path=_project_relative(path, project_root),
            )
            _registry[attachment_id] = record
            registered.append({'ok': True, **_record_to_dict(record)})

    return registered


def release_attachments(ids: Iterable[str]) -> int:
    removed = 0
    with _lock:
        for attachment_id in ids:
            key = str(attachment_id)
            if _registry.pop(key, None) is not None:
                removed += 1
            _content_cache.pop(key, None)
    return removed


def attachment_metadata(ids: Iterable[str]) -> list[dict]:
    rows: list[dict] = []
    with _lock:
        _cleanup_locked()
        for attachment_id in ids:
            record = _registry.get(str(attachment_id))
            if record:
                rows.append(_record_to_dict(record))
    return rows


def _decode_text_bytes(raw: bytes) -> str:
    if not raw:
        return ''
    if b'\x00' in raw[:4096]:
        raise ValueError('binary')
    last_error: Exception | None = None
    for encoding in ('utf-8-sig', 'utf-8', 'cp949', 'euc-kr'):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError as exc:
            last_error = exc
    raise ValueError(f'text-decode-failed: {last_error}')


def _read_plain_text(path: Path) -> str:
    with path.open('rb') as stream:
        raw = stream.read(_MAX_FILE_BYTES + 1)
    if len(raw) > _MAX_FILE_BYTES:
        raise ValueError('파일 크기 제한 초과')
    return _decode_text_bytes(raw)


def _read_notebook(path: Path) -> str:
    data = json.loads(_read_plain_text(path))
    cells = data.get('cells') if isinstance(data, dict) else []
    rows: list[str] = []
    for index, cell in enumerate(cells or []):
        if not isinstance(cell, dict):
            continue
        cell_type = str(cell.get('cell_type') or 'unknown')
        source = cell.get('source') or []
        if isinstance(source, list):
            source = ''.join(str(value) for value in source)
        rows.append(f'## Cell {index + 1} ({cell_type})\n{source}')
    return '\n\n'.join(rows)


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise ValueError('PDF 분석 모듈(pypdf)이 설치되어 있지 않습니다.') from exc
    reader = PdfReader(str(path))
    chunks: list[str] = []
    for index, page in enumerate(reader.pages[:80]):
        text = str(page.extract_text() or '').strip()
        if text:
            chunks.append(f'## Page {index + 1}\n{text}')
        if sum(len(x) for x in chunks) >= _MAX_FILE_CONTEXT_CHARS * 2:
            break
    return '\n\n'.join(chunks)


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except Exception as exc:
        raise ValueError('Word 분석 모듈(python-docx)이 설치되어 있지 않습니다.') from exc
    doc = Document(str(path))
    rows = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                rows.append(' | '.join(values))
    return '\n'.join(rows)


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        raise ValueError('Excel 분석 모듈(openpyxl)이 설치되어 있지 않습니다.') from exc
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[str] = []
    try:
        for sheet in workbook.worksheets[:20]:
            chunks.append(f'## Sheet: {sheet.title}')
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                values = ['' if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    chunks.append('\t'.join(values))
                if row_index >= 500:
                    chunks.append('[이후 행 생략]')
                    break
                if sum(len(x) for x in chunks) >= _MAX_FILE_CONTEXT_CHARS * 2:
                    break
            if sum(len(x) for x in chunks) >= _MAX_FILE_CONTEXT_CHARS * 2:
                break
    finally:
        workbook.close()
    return '\n'.join(chunks)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ValueError('PowerPoint 분석 모듈(python-pptx)이 설치되어 있지 않습니다.') from exc
    presentation = Presentation(str(path))
    chunks: list[str] = []
    for index, slide in enumerate(presentation.slides[:100]):
        lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, 'text', '')
            if str(text).strip():
                lines.append(str(text).strip())
        if lines:
            chunks.append(f'## Slide {index + 1}\n' + '\n'.join(lines))
        if sum(len(x) for x in chunks) >= _MAX_FILE_CONTEXT_CHARS * 2:
            break
    return '\n\n'.join(chunks)


def _extract_content(record: AttachmentRecord) -> tuple[str, str]:
    path = Path(record.path)
    if not path.exists() or not path.is_file():
        raise ValueError('등록 후 파일이 이동되었거나 삭제되었습니다.')
    if path.stat().st_size > _MAX_FILE_BYTES:
        raise ValueError('등록 후 파일 크기가 제한을 초과했습니다.')

    suffix = path.suffix.casefold()
    if suffix == '.ipynb':
        return _read_notebook(path), 'notebook'
    if suffix == '.pdf':
        return _read_pdf(path), 'pdf'
    if suffix == '.docx':
        return _read_docx(path), 'docx'
    if suffix in {'.xlsx', '.xlsm'}:
        return _read_xlsx(path), 'xlsx'
    if suffix == '.pptx':
        return _read_pptx(path), 'pptx'
    if suffix in _TEXT_EXTENSIONS or not suffix:
        return _read_plain_text(path), 'text'

    # Some extension-less or uncommon source files are still text.  Probe them
    # instead of rejecting by extension alone.
    try:
        return _read_plain_text(path), 'text'
    except Exception as exc:
        raise ValueError(f'현재 AI 첨부 분석에서 지원하지 않는 바이너리 형식입니다: {suffix or path.name}') from exc


def _file_signature(path: Path) -> tuple[int, int]:
    stat = path.stat()
    return int(stat.st_size), int(getattr(stat, 'st_mtime_ns', int(stat.st_mtime * 1_000_000_000)))


def prepare_attachment(attachment_id: str) -> dict:
    """Extract one registered attachment and cache the extracted text.

    The returned payload intentionally contains metadata only. The extracted
    file contents stay inside the Backend and are reused by
    ``build_attachment_context`` so progress reporting never sends document
    bodies through the progress channel.
    """
    key = str(attachment_id or '').strip()
    if not key:
        raise ValueError('첨부 ID가 없습니다.')

    with _lock:
        _cleanup_locked()
        record = _registry.get(key)
    if not record:
        raise ValueError(f'첨부 ID가 만료되었거나 존재하지 않습니다: {key[:8]}')

    path = Path(record.path)
    if not path.exists() or not path.is_file():
        raise ValueError('등록 후 파일이 이동되었거나 삭제되었습니다.')
    size, mtime_ns = _file_signature(path)
    if size > _MAX_FILE_BYTES:
        raise ValueError('등록 후 파일 크기가 제한을 초과했습니다.')

    with _lock:
        cached = _content_cache.get(key)
        if cached and cached.get('size') == size and cached.get('mtime_ns') == mtime_ns:
            content = str(cached.get('content') or '')
            return {
                **_record_to_dict(record),
                'content_type': str(cached.get('content_type') or 'text'),
                'content_chars': len(content),
                'cached': True,
            }

    content, content_type = _extract_content(record)
    content = str(content or '')
    with _lock:
        # The file could have changed while a document parser was running. If
        # it did, do not cache stale extracted text; the next request retries.
        current_path = Path(record.path)
        current_size, current_mtime_ns = _file_signature(current_path)
        if current_size == size and current_mtime_ns == mtime_ns:
            _content_cache[key] = {
                'content': content,
                'content_type': content_type,
                'size': size,
                'mtime_ns': mtime_ns,
                'extracted_at': time.time(),
            }
    return {
        **_record_to_dict(record),
        'content_type': content_type,
        'content_chars': len(content),
        'cached': False,
    }


def _prepared_content(record: AttachmentRecord) -> tuple[str, str]:
    prepared = prepare_attachment(record.attachment_id)
    with _lock:
        cached = _content_cache.get(record.attachment_id)
        if cached:
            return str(cached.get('content') or ''), str(cached.get('content_type') or prepared.get('content_type') or 'text')
    # A file that changed during extraction is deliberately not cached. Read it
    # once more so the actual AI request remains authoritative.
    return _extract_content(record)


def build_attachment_context(
    ids: Iterable[str],
    *,
    purpose: str = 'AI 분석',
    total_char_limit: int = _DEFAULT_TOTAL_CONTEXT_CHARS,
    per_file_char_limit: int = _MAX_FILE_CONTEXT_CHARS,
) -> dict:
    requested = [str(value) for value in ids if str(value or '').strip()]
    if not requested:
        return {'text': '', 'files': [], 'warnings': [], 'total_chars': 0}

    files: list[dict] = []
    warnings: list[str] = []
    sections: list[str] = []
    total_chars = 0

    with _lock:
        _cleanup_locked()
        records = [(_registry.get(value), value) for value in requested]

    for record, requested_id in records:
        if not record:
            warnings.append(f'첨부 ID가 만료되었거나 존재하지 않습니다: {requested_id[:8]}')
            continue
        try:
            content, content_type = _prepared_content(record)
        except Exception as exc:
            warnings.append(f'{record.name}: {exc}')
            files.append({**_record_to_dict(record), 'included': False, 'message': str(exc)})
            continue

        content = redact_sensitive_text(str(content or '')).strip()
        if not content:
            warnings.append(f'{record.name}: 분석할 텍스트를 추출하지 못했습니다.')
            files.append({**_record_to_dict(record), 'included': False, 'message': '텍스트 없음'})
            continue

        remaining = max(0, int(total_char_limit) - total_chars)
        if remaining <= 0:
            warnings.append('첨부 파일 전체 Context 예산을 초과하여 일부 파일을 생략했습니다.')
            break

        allowed = min(max(1000, int(per_file_char_limit)), remaining)
        clipped = content[:allowed]
        truncated = len(content) > len(clipped)
        label = record.project_relative_path or record.path
        section = (
            f'### 첨부 파일: {record.name}\n'
            f'- 경로: {label}\n'
            f'- 형식: {content_type}\n'
            f'- 용도: {purpose}\n\n'
            f'{clipped}'
        )
        if truncated:
            section += '\n\n[파일 내용이 Context 예산에 맞게 일부 생략됨]'
        sections.append(section)
        total_chars += len(section)
        files.append({
            **_record_to_dict(record),
            'included': True,
            'content_type': content_type,
            'context_chars': len(clipped),
            'truncated': truncated,
        })

    text = ''
    if sections:
        text = '[사용자가 직접 선택하여 등록한 참고 파일]\n' + '\n\n'.join(sections)
    return {
        'text': text,
        'files': files,
        'warnings': warnings,
        'total_chars': total_chars,
    }


_REQUIREMENTS_TOTAL_CONTEXT_CHARS = 28_000
_REQUIREMENTS_PER_FILE_CHARS = 12_000


def _unique_lines(lines: Iterable[str], limit: int = 160) -> list[str]:
    result: list[str] = []
    seen: set[str] = set()
    for raw in lines:
        line = str(raw or '').strip()
        if not line:
            continue
        key = re.sub(r'\s+', ' ', line).casefold()
        if key in seen:
            continue
        seen.add(key)
        result.append(line)
        if len(result) >= limit:
            break
    return result


def _requirement_candidate_lines(content: str, limit: int = 220) -> list[str]:
    """Extract requirement-bearing natural language before source-code compression.

    Notebook problem statements are commonly written in Markdown cells.  Older
    requirements digests kept imports/framework clues but could drop ordinary
    Korean bullets and paragraphs, which meant a requirement-rich assignment was
    reduced to only a handful of technologies.  This pass keeps Markdown prose,
    bullets, constraints and deliverables as first-class evidence.
    """
    text = str(content or '').replace('\r\n', '\n').replace('\r', '\n')
    lines = text.splitlines()
    result: list[str] = []
    seen: set[str] = set()
    notebook_cell = ''
    notebook_cell_type = ''

    requirement_words = re.compile(
        r'(?i)(해야|한다|하도록|필수|반드시|금지|제한|사용|구현|제공|저장|조회|검색|처리|관리|구성|연결|'
        r'목표|산출물|제약|제공 자료|환경 확인|streamlit|fastapi|postgres|pgvector|redis|langchain|langgraph|'
        r'openai|ollama|mcp|csv|api key|password|환경변수|최대\s*\d+|\d+\s*건)'
    )
    bullet_re = re.compile(r'^\s*(?:[-*+]\s+|\d+[.)]\s+)')
    cell_re = re.compile(r'^##\s+Cell\s+(\d+)\s+\(([^)]+)\)\s*$', re.IGNORECASE)

    def add(raw: str, location: str = '') -> None:
        value = str(raw or '').strip()
        if not value or value.startswith('```'):
            return
        if len(value) > 900:
            value = value[:900]
        # Skip obvious source-code statements. Comments in code cells may still be
        # useful and are handled by the structural clue pass below.
        if value.startswith(('import ', 'from ', 'def ', 'class ', 'async def ', 'SELECT ', 'CREATE ', 'INSERT ', 'UPDATE ')):
            return
        cleaned = re.sub(r'^\s*(?:[-*+]\s+|\d+[.)]\s+)', '', value).strip()
        cleaned = re.sub(r'^#{1,6}\s+', '', cleaned).strip()
        if not cleaned:
            return
        key = re.sub(r'\s+', ' ', cleaned).casefold()
        if key in seen:
            return
        seen.add(key)
        prefix = f'[{location}] ' if location else ''
        result.append(prefix + cleaned)

    for line in lines:
        match = cell_re.match(str(line or '').strip())
        if match:
            notebook_cell = f'Cell {match.group(1)} {match.group(2).strip()}'
            notebook_cell_type = match.group(2).strip().casefold()
            continue

        stripped = str(line or '').strip()
        if not stripped:
            continue

        is_markdown = notebook_cell_type == 'markdown'
        is_bullet = bool(bullet_re.match(stripped))
        is_heading = bool(re.match(r'^#{1,6}\s+', stripped))
        has_requirement_signal = bool(requirement_words.search(stripped))
        has_korean_prose = bool(re.search(r'[가-힣]', stripped)) and not re.search(r'[{}();=]{3,}', stripped)

        if is_markdown and (is_bullet or is_heading or has_requirement_signal or has_korean_prose):
            # Long assignment paragraphs frequently contain several independent
            # requirements. Split on sentence boundaries before adding them.
            pieces = re.split(r'(?<=[.!?。])\s+', stripped)
            for piece in pieces:
                add(piece, notebook_cell)
        elif is_bullet and (has_requirement_signal or has_korean_prose):
            add(stripped, notebook_cell)
        elif has_requirement_signal and has_korean_prose:
            add(stripped, notebook_cell)

        if len(result) >= limit:
            break

    return result[:limit]


def _requirements_outline(content: str, content_type: str, limit: int) -> str:
    """Build a requirement-first, bounded evidence digest for design interviews.

    Explicit Markdown requirements are preserved before framework/import clues.
    This prevents large Notebook assignments from losing ordinary prose bullets
    when a local LLM receives the compact context.
    """
    text = str(content or '').replace('\r\n', '\n').replace('\r', '\n').strip()
    if not text:
        return ''

    target = max(3_000, int(limit))
    candidate_lines = _requirement_candidate_lines(text, limit=240)
    candidate_block = '\n'.join(f'- {line}' for line in candidate_lines)

    lines = text.splitlines()
    selected: list[str] = []
    patterns = (
        r'^#{1,6}\s+',
        r'^##\s+Cell\s+\d+',
        r'^(?:class|def|async\s+def)\s+',
        r'^(?:from\s+\S+\s+import|import\s+\S+)',
        r'^(?:CREATE|ALTER|INSERT|UPDATE|DELETE|SELECT)\b',
        r'\b(?:streamlit|st\.|fastapi|uvicorn|redis|psycopg|postgres|pgvector|openai|ollama|langchain|langgraph|mcp|supabase|firestore|embedding|vector|rag|router|endpoint|websocket|sse)\b',
        r'^\s*(?:#|//|--|/\*|\*)\s*[^\s]',
    )
    combined = re.compile('|'.join(f'(?:{pat})' for pat in patterns), re.IGNORECASE)
    for line in lines:
        if combined.search(line):
            selected.append(line)
    selected = _unique_lines(selected, 220)
    structure = '\n'.join(selected)

    # Requirements receive most of the budget.  Structure clues are secondary,
    # and small head/tail excerpts remain only as a safety net.
    candidate_budget = min(len(candidate_block), max(1_800, int(target * 0.62)))
    structure_budget = max(900, int(target * 0.22))
    excerpt_budget = max(600, target - candidate_budget - structure_budget - 350)
    head_budget = max(400, int(excerpt_budget * 0.7))
    tail_budget = max(200, excerpt_budget - head_budget)

    chunks: list[str] = []
    if candidate_block:
        chunks.append('[명시적 요구사항 후보]\n' + candidate_block[:candidate_budget])
    if structure:
        chunks.append('[기술/구조 단서]\n' + structure[:structure_budget])
    if not candidate_block or len(candidate_block) < candidate_budget // 2:
        chunks.append('[대표 앞부분]\n' + text[:head_budget])
        if len(text) > head_budget + tail_budget:
            chunks.append('[대표 끝부분]\n' + text[-tail_budget:])

    compact = '\n\n'.join(part for part in chunks if str(part).strip())
    return compact[:target]

def build_requirements_attachment_context(
    ids: Iterable[str],
    *,
    purpose: str = 'Agent 설계 인터뷰 요구사항/참고자료 분석',
    total_char_limit: int = _REQUIREMENTS_TOTAL_CONTEXT_CHARS,
    per_file_char_limit: int = _REQUIREMENTS_PER_FILE_CHARS,
) -> dict:
    """Return a compact attachment context specialized for requirements interviews.

    Raw attachment bodies remain cached inside the Backend.  Only a bounded
    structural/evidence digest is injected into the interview LLM, preventing
    local models from losing the system instruction or echoing large source files.
    """
    requested = [str(value) for value in ids if str(value or '').strip()]
    if not requested:
        return {'text': '', 'files': [], 'warnings': [], 'total_chars': 0, 'mode': 'requirements_digest'}

    files: list[dict] = []
    warnings: list[str] = []
    sections: list[str] = []
    total_chars = 0

    with _lock:
        _cleanup_locked()
        records = [(_registry.get(value), value) for value in requested]

    # Share the requirements context budget across all selected files.  Older
    # behavior allowed the first two large files to consume the whole budget and
    # silently omit a later primary Notebook problem statement.
    fair_share = min(
        max(1_500, int(per_file_char_limit)),
        max(3_000, int(total_char_limit) // max(1, len(records))),
    )

    for record, requested_id in records:
        if not record:
            warnings.append(f'첨부 ID가 만료되었거나 존재하지 않습니다: {requested_id[:8]}')
            continue
        try:
            content, content_type = _prepared_content(record)
        except Exception as exc:
            warnings.append(f'{record.name}: {exc}')
            files.append({**_record_to_dict(record), 'included': False, 'message': str(exc)})
            continue

        content = redact_sensitive_text(str(content or '')).strip()
        if not content:
            warnings.append(f'{record.name}: 분석할 텍스트를 추출하지 못했습니다.')
            files.append({**_record_to_dict(record), 'included': False, 'message': '텍스트 없음'})
            continue

        remaining = max(0, int(total_char_limit) - total_chars)
        if remaining <= 0:
            warnings.append('Agent 설계 인터뷰용 첨부 Context 예산을 초과하여 일부 파일을 생략했습니다.')
            break

        allowed = min(fair_share, remaining)
        digest = _requirements_outline(content, content_type, allowed).strip()
        if not digest:
            warnings.append(f'{record.name}: 인터뷰용 분석 요약을 만들지 못했습니다.')
            files.append({**_record_to_dict(record), 'included': False, 'message': '요약 없음'})
            continue

        label = record.project_relative_path or record.path
        section = (
            f'### 참고 파일 분석본: {record.name}\n'
            f'- 경로: {label}\n'
            f'- 형식: {content_type}\n'
            f'- 원문 문자 수: {len(content)}\n'
            f'- 용도: {purpose}\n'
            f'- 주의: 아래 내용은 요구사항 파악을 위한 압축 분석본입니다. 원문을 그대로 답변에 복사하지 마세요.\n\n'
            f'{digest}'
        )
        sections.append(section)
        total_chars += len(section)
        files.append({
            **_record_to_dict(record),
            'included': True,
            'content_type': content_type,
            'source_chars': len(content),
            'context_chars': len(digest),
            'truncated': len(content) > len(digest),
            'context_mode': 'requirements_digest',
        })

    text = ''
    if sections:
        text = (
            '[사용자가 직접 선택한 참고 파일의 인터뷰용 압축 분석 Context]\n'
            '중요: 이 Context는 내부 분석 근거입니다. 답변에 파일 원문/코드 블록을 대량 출력하지 말고, '
            '확인한 요구사항을 짧게 반영한 뒤 다음 질문 하나만 하세요.\n\n'
            + '\n\n'.join(sections)
        )
    return {
        'text': text,
        'files': files,
        'warnings': warnings,
        'total_chars': total_chars,
        'mode': 'requirements_digest',
    }
