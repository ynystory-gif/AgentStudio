from __future__ import annotations

from io import BytesIO
from pathlib import Path
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "backend"))

APP = (ROOT / "frontend" / "src" / "App.jsx").read_text(encoding="utf-8")
MAIN = (ROOT / "backend" / "app" / "main.py").read_text(encoding="utf-8")
ROUTES = (ROOT / "backend" / "app" / "api" / "routes.py").read_text(encoding="utf-8")
SERVICE = (ROOT / "backend" / "app" / "services" / "presentation_export_service.py").read_text(encoding="utf-8")
ICON_DIR = ROOT / "backend" / "app" / "services" / "ppt_assets" / "large_icons"


def require(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


require("AGENTSTUDIO_FRONTEND_VERSION='5.369'" in APP, "frontend version must be 5.369")
require('version="5.369"' in MAIN or "version='5.369'" in MAIN, "backend version must be 5.369")
require('"version": "5.369"' in ROUTES, "health version must be 5.369")
require("_ICON_RULES" in SERVICE and "_add_large_visual" in SERVICE, "large visual icon registry missing")
require("_add_project_stack_slide" in SERVICE, "project adaptive stack visual slide missing")

required_icons = {
    "user_group", "laptop", "mobile", "globe", "shield", "api_gateway",
    "agent_cube", "workflow", "llm", "mcp", "tool", "database", "cache",
    "vector_db", "storage", "cloud", "server", "kubernetes", "network",
    "report", "terminal", "code", "success", "warning", "error", "generic_component",
}
missing = sorted(name for name in required_icons if not (ICON_DIR / f"{name}.png").is_file())
require(not missing, f"missing large visual assets: {missing}")

from app.services.presentation_export_service import build_agentstudio_presentation, _icon_key_for_text

require(_icon_key_for_text("MCP Tool Registry", "MCP Client / Tool") == "mcp", "MCP must map to MCP visual")
require(_icon_key_for_text("Redis Cache") == "cache", "Redis must map to cache visual")
require(_icon_key_for_text("PostgreSQL") == "database", "PostgreSQL must map to database visual")
require(_icon_key_for_text("FastAPI REST API") == "api_gateway", "FastAPI API must map to API visual")

payload = {
    "scope": "ALL",
    "project_name": "v5353_visual_contract",
    "workflow_definition": {},
    "coding_style_report": {"fail": 0, "warning": 0},
    "report": {
        "status": "COMPLETED",
        "testReturncode": 0,
        "targetWorkflow": {"name": "Agent Workflow", "steps": [{"label": "Input"}, {"label": "Agent"}, {"label": "MCP Tool"}]},
        "requirementSpec": {"goal": "Large visual architecture export"},
        "architecture": {
            "interfaces": [{"name": "FastAPI REST API"}],
            "components": [
                {"name": "Input Router", "description": "Input"},
                {"name": "Agent Orchestrator", "description": "LangGraph Agent"},
                {"name": "OpenAI LLM", "description": "LLM"},
                {"name": "MCP Tool Registry", "description": "MCP Client / Tool"},
                {"name": "Response Renderer", "description": "Output"},
            ],
            "persistence": [{"name": "PostgreSQL"}, {"name": "Redis Cache"}, {"name": "pgvector"}, {"name": "File Storage"}, {"name": "Report State"}],
            "security": [{"name": "OAuth"}],
            "state": [{"name": "session"}],
        },
        "asBuiltArchitecture": {"components": [], "interfaces": [], "persistence": [], "scan": {"source_file_count": 0}},
        "architectureConformance": {},
        "projectProfile": {
            "project_type": "AI_AGENT",
            "project_type_label": "AI Agent / LLM Application",
            "tech_stack": ["FastAPI", "LangGraph", "OpenAI", "MCP", "PostgreSQL", "Redis"],
        },
    },
}
content, _ = build_agentstudio_presentation(payload, "5.369")
prs = Presentation(BytesIO(content))
require(len(prs.slides) == 6, f"Agent ALL export must contain 6 target-only adaptive slides, got {len(prs.slides)}")
picture_counts = [sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE) for slide in prs.slides]
require(picture_counts[4] >= 10, f"design architecture must contain large visual pictures, got {picture_counts[4]}")
require(picture_counts[5] >= 8, f"project technology/runtime slide must contain large visual pictures, got {picture_counts[5]}")

print("PASS v5.369 Large Architecture Visual Asset contract")
