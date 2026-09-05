from __future__ import annotations

import csv
import io
import json
import re
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path
from typing import Iterable

from app.rag.constants import SENSITIVE_PARTS, SOURCE_CODE_EXTENSIONS


MAX_FILE_BYTES = 20 * 1024 * 1024
MAX_TABULAR_ROWS = 50_000

TEXT_EXTENSIONS = {
    '.txt', '.md', '.markdown', '.html', '.htm', '.csv', '.json', '.xml', '.yaml', '.yml',
    '.toml', '.ini', '.cfg', '.conf', '.log', '.rst', '.sql', '.py', '.pyw', '.ts', '.tsx',
    '.js', '.jsx', '.java', '.cs', '.go', '.rs', '.cpp', '.c', '.h', '.hpp', '.ps1', '.sh',
    '.bat', '.cmd', '.vue', '.svelte', '.gradle', '.properties', '.env.example',
}
DOCUMENT_EXTENSIONS = TEXT_EXTENSIONS | {'.pdf', '.docx', '.pptx', '.xlsx'}

LANGUAGE_BY_SUFFIX = {
    '.py': 'Python', '.pyw': 'Python', '.ts': 'TypeScript', '.tsx': 'TypeScript React',
    '.js': 'JavaScript', '.jsx': 'JavaScript React', '.java': 'Java', '.cs': 'C#', '.go': 'Go',
    '.rs': 'Rust', '.cpp': 'C++', '.c': 'C', '.h': 'C/C++ Header', '.hpp': 'C++ Header',
    '.sql': 'SQL', '.ps1': 'PowerShell', '.sh': 'Shell', '.bat': 'Batch', '.cmd': 'Batch',
    '.yaml': 'YAML', '.yml': 'YAML', '.json': 'JSON', '.xml': 'XML', '.toml': 'TOML',
    '.md': 'Markdown', '.markdown': 'Markdown', '.html': 'HTML', '.htm': 'HTML',
}


@dataclass(slots=True)
class LoadedDocument:
    path: Path
    display_path: str
    filename: str
    document_type: str
    language: str
    text: str
    size_bytes: int
    metadata: dict = field(default_factory=dict)


class _VisibleTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._ignored = 0
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:  # noqa: ANN001
        if tag.lower() in {'script', 'style', 'noscript'}:
            self._ignored += 1
        elif tag.lower() in {'p', 'div', 'section', 'article', 'li', 'br', 'tr', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'}:
            self.parts.append('\n')

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in {'script', 'style', 'noscript'}:
            self._ignored = max(0, self._ignored - 1)
        elif tag.lower() in {'p', 'div', 'section', 'article', 'li', 'tr', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'}:
            self.parts.append('\n')

    def handle_data(self, data: str) -> None:
        if not self._ignored and data.strip():
            self.parts.append(data)


def resolve_source_path(project_root: str, source_uri: str) -> Path:
    raw = str(source_uri or '').strip()
    if not raw:
        raise ValueError('Source 경로를 입력하세요.')
    root = Path(project_root).expanduser().resolve() if project_root else None
    path = Path(raw).expanduser()
    if not path.is_absolute() and root is not None:
        path = root / path
    resolved = path.resolve()
    if root is not None:
        try:
            resolved.relative_to(root)
        except ValueError as exc:
            raise ValueError('현재 Agent 프로젝트 밖의 Source는 RAG Indexing에서 허용하지 않습니다.') from exc
    if not resolved.exists():
        raise FileNotFoundError(f'Source를 찾을 수 없습니다: {resolved}')
    return resolved


def detect_document_type(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    language = LANGUAGE_BY_SUFFIX.get(suffix, '')
    if suffix == '.sql':
        return 'SQL', 'SQL'
    if suffix in SOURCE_CODE_EXTENSIONS and suffix not in {'.md', '.json', '.xml', '.yaml', '.yml'}:
        return 'SOURCE_CODE', language or 'Source Code'
    if suffix in {'.md', '.markdown'}:
        return 'MARKDOWN', 'Markdown'
    if suffix == '.pdf':
        return 'PDF', ''
    if suffix == '.docx':
        return 'DOCX', ''
    if suffix == '.pptx':
        return 'PPTX', ''
    if suffix == '.xlsx':
        return 'EXCEL', ''
    if suffix == '.csv':
        return 'CSV', ''
    if suffix in {'.html', '.htm'}:
        return 'HTML', 'HTML'
    if suffix == '.json':
        return 'JSON', 'JSON'
    if suffix == '.xml':
        return 'XML', 'XML'
    if suffix in {'.yaml', '.yml'}:
        return 'YAML', 'YAML'
    if suffix == '.sql':
        return 'SQL', 'SQL'
    return 'TEXT', language


def _decode_bytes(data: bytes) -> str:
    for encoding in ('utf-8-sig', 'utf-8', 'cp949', 'euc-kr'):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode('latin-1', errors='replace')


def _read_plain_text(path: Path) -> str:
    return _decode_bytes(path.read_bytes())


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = str(page.extract_text() or '').strip()
        if text:
            pages.append(f'# Page {page_index}\n{text}')
    return '\n\n'.join(pages)


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        value = paragraph.text.strip()
        if value:
            parts.append(value)
    for table_index, table in enumerate(doc.tables, start=1):
        parts.append(f'\n# Table {table_index}')
        for row in table.rows:
            values = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            if any(values):
                parts.append('\t'.join(values))
    return '\n'.join(parts)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    deck = Presentation(str(path))
    parts: list[str] = []
    for slide_index, slide in enumerate(deck.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = str(getattr(shape, 'text', '') or '').strip()
            if text:
                lines.append(text)
        if lines:
            parts.append(f'# Slide {slide_index}\n' + '\n'.join(lines))
    return '\n\n'.join(parts)


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    seen_rows = 0
    try:
        for sheet in workbook.worksheets:
            parts.append(f'# Sheet: {sheet.title}')
            for row in sheet.iter_rows(values_only=True):
                values = ['' if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    parts.append('\t'.join(values))
                    seen_rows += 1
                if seen_rows >= MAX_TABULAR_ROWS:
                    parts.append('[행 수 제한으로 이후 데이터 생략]')
                    return '\n'.join(parts)
    finally:
        workbook.close()
    return '\n'.join(parts)


def _read_csv(path: Path) -> str:
    raw = _decode_bytes(path.read_bytes())
    reader = csv.reader(io.StringIO(raw))
    parts: list[str] = []
    for index, row in enumerate(reader):
        if index >= MAX_TABULAR_ROWS:
            parts.append('[행 수 제한으로 이후 데이터 생략]')
            break
        parts.append('\t'.join(str(value) for value in row))
    return '\n'.join(parts)


def _read_html(path: Path) -> str:
    parser = _VisibleTextParser()
    parser.feed(_read_plain_text(path))
    return re.sub(r'\n{3,}', '\n\n', ''.join(parser.parts)).strip()


def _read_json(path: Path) -> str:
    raw = _read_plain_text(path)
    try:
        payload = json.loads(raw)
        return json.dumps(payload, ensure_ascii=False, indent=2)
    except Exception:
        return raw


def read_document_text(path: Path, document_type: str) -> str:
    if document_type == 'PDF':
        return _read_pdf(path)
    if document_type == 'DOCX':
        return _read_docx(path)
    if document_type == 'PPTX':
        return _read_pptx(path)
    if document_type == 'EXCEL':
        return _read_xlsx(path)
    if document_type == 'CSV':
        return _read_csv(path)
    if document_type == 'HTML':
        return _read_html(path)
    if document_type == 'JSON':
        return _read_json(path)
    return _read_plain_text(path)


def _display_path(path: Path, project_root: str) -> str:
    if project_root:
        try:
            return str(path.resolve().relative_to(Path(project_root).expanduser().resolve())).replace('\\', '/')
        except Exception:
            pass
    return str(path).replace('\\', '/')


def _supported_file(path: Path, source_type: str) -> bool:
    suffix = path.suffix.lower()
    if source_type == 'SOURCE_CODE':
        return suffix in SOURCE_CODE_EXTENSIONS
    return suffix in DOCUMENT_EXTENSIONS or suffix in SOURCE_CODE_EXTENSIONS


def iter_source_files(project_root: str, source_type: str, source_uri: str) -> Iterable[Path]:
    path = resolve_source_path(project_root, source_uri)
    source_type = str(source_type or 'FILE').upper()
    if path.is_file():
        yield path
        return
    if not path.is_dir():
        raise ValueError('RAG Source는 File 또는 Folder여야 합니다.')

    for candidate in path.rglob('*'):
        if not candidate.is_file():
            continue
        lower_parts = {part.lower() for part in candidate.parts}
        if lower_parts & SENSITIVE_PARTS:
            continue
        if _supported_file(candidate, source_type):
            yield candidate


def load_source_documents(project_root: str, source_type: str, source_uri: str) -> tuple[list[LoadedDocument], list[dict]]:
    documents: list[LoadedDocument] = []
    skipped: list[dict] = []
    for path in iter_source_files(project_root, source_type, source_uri):
        try:
            if not _supported_file(path, str(source_type or 'FILE').upper()):
                skipped.append({'path': _display_path(path, project_root), 'reason': f'2차 Indexing 미지원 파일 형식: {path.suffix or "확장자 없음"}'})
                continue
            size_bytes = int(path.stat().st_size)
            if size_bytes > MAX_FILE_BYTES:
                skipped.append({'path': _display_path(path, project_root), 'reason': f'파일 크기 제한 초과 ({size_bytes} bytes)'})
                continue
            document_type, language = detect_document_type(path)
            text = read_document_text(path, document_type).replace('\x00', '').strip()
            if not text:
                skipped.append({'path': _display_path(path, project_root), 'reason': '추출 가능한 텍스트가 없습니다.'})
                continue
            documents.append(LoadedDocument(
                path=path,
                display_path=_display_path(path, project_root),
                filename=path.name,
                document_type=document_type,
                language=language,
                text=text,
                size_bytes=size_bytes,
                metadata={'suffix': path.suffix.lower()},
            ))
        except Exception as exc:
            skipped.append({'path': _display_path(path, project_root), 'reason': str(exc)})
    return documents, skipped
