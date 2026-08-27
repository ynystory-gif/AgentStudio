from __future__ import annotations

from datetime import datetime
from io import BytesIO
from pathlib import Path
import json
import re
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# 16:9 wide layout
SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)
FONT = "Malgun Gothic"
FONT_MONO = "Consolas"

# Light presentation palette inspired by the application's blue/purple UI,
# intentionally implemented with native PowerPoint objects so users can edit it.
C = {
    "navy": "17365D",
    "navy2": "244A7C",
    "blue": "2F6BFF",
    "blue_soft": "EAF2FF",
    "cyan": "28BFC7",
    "cyan_soft": "E8FBFC",
    "purple": "7657E8",
    "purple_soft": "F1EDFF",
    "green": "1C9C69",
    "green_soft": "EAF8F2",
    "orange": "E08A28",
    "orange_soft": "FFF4E5",
    "red": "D94B55",
    "red_soft": "FDEDEF",
    "ink": "162033",
    "muted": "64748B",
    "line": "D6E1F0",
    "panel": "F7FAFE",
    "white": "FFFFFF",
    "black": "000000",
}

_SCOPE_LABEL = {
    "ALL": "전체",
    "WORKFLOW": "워크플로우",
    "RUN": "실행결과",
    "REPORT": "분석리포트",
    "ARCHITECTURE": "아키텍처",
    "DB_ERD": "DB_ERD",
}


# v5.363 Large Architecture Visual Asset Pack
# These are intentionally generic architecture illustrations (not emoji/vendor logos).
# Each PNG is a separate PowerPoint picture object, while labels/cards/connectors remain editable native shapes.
_ICON_DIR = Path(__file__).resolve().parent / "ppt_assets" / "large_icons"

_ICON_RULES: list[tuple[tuple[str, ...], str]] = [
    (("mcp",), "mcp"),
    (("user", "client", "customer", "사용자", "고객"), "user_group"),
    (("mobile", "phone", "smartphone", "모바일", "스마트폰"), "mobile"),
    (("react", "vite", "frontend", "web ui", "browser", "desktop", "workspace", "프론트", "웹"), "laptop"),
    (("interface", "api", "fastapi", "gateway", "endpoint", "인터페이스"), "api_gateway"),
    (("security", "auth", "oauth", "jwt", "permission", "보안", "인증"), "shield"),
    (("orchestrator", "agent", "에이전트"), "agent_cube"),
    (("workflow", "langgraph", "planner", "planning", "state", "flow", "워크플로우", "계획"), "workflow"),
    (("openai", "ollama", "llm", "langchain", "model", "모델"), "llm"),
    (("tool", "plugin", "function", "도구"), "tool"),
    (("redis", "cache", "캐시"), "cache"),
    (("pgvector", "vector", "embedding", "벡터"), "vector_db"),
    (("postgres", "postgresql", "mssql", "oracle", "sqlite", "sql", "database", "db", "데이터베이스"), "database"),
    (("bucket", "storage", "file", "artifact", "파일", "스토리지"), "storage"),
    (("kubernetes", "k8s", "cluster", "쿠버네티스"), "kubernetes"),
    (("cloud", "supabase", "firestore", "클라우드"), "cloud"),
    (("docker", "container", "server", "runtime", "backend", "서버"), "server"),
    (("network", "socket", "websocket", "sse", "네트워크"), "network"),
    (("report", "analysis", "log", "리포트", "분석", "로그"), "report"),
    (("terminal", "console", "shell", "터미널"), "terminal"),
    (("python", "code", "source", "코드"), "code"),
]


def _icon_key_for_text(*values: Any, fallback: str = "generic_component") -> str:
    text = " ".join(_safe_text(v, "", 300).lower() for v in values if v is not None)
    for needles, key in _ICON_RULES:
        if any(needle in text for needle in needles):
            return key
    return fallback


def _icon_path(key: str) -> Path | None:
    requested = _ICON_DIR / f"{key}.png"
    if requested.is_file():
        return requested
    generic = _ICON_DIR / "generic_component.png"
    return generic if generic.is_file() else None


def _add_large_visual(slide, key: str, x: float, y: float, w: float, h: float):
    path = _icon_path(key)
    if path is not None:
        try:
            return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        except Exception:
            pass
    # Native-shape fallback: PowerPoint export must never fail because an asset is missing.
    fallback = _add_box(slide, x, y, w, h, fill=C["blue_soft"], line=C["blue"], radius=True)
    _set_shape_text(fallback, "ARCH", "visual", title_size=max(8, min(16, w * 8)), subtitle_size=6.5, title_color=C["blue"])
    return fallback


def _add_visual_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    title: str,
    subtitle: str = "",
    icon_key: str = "generic_component",
    fill: str = C["white"],
    line: str = C["line"],
    icon_ratio: float = 0.53,
    title_size: float = 10.0,
    subtitle_size: float = 7.0,
):
    card = _add_box(slide, x, y, w, h, fill=fill, line=line)
    icon_h = min(h * icon_ratio, w * 0.72)
    icon_w = min(w * 0.64, icon_h * 1.18)
    icon_x = x + (w - icon_w) / 2
    icon_y = y + 0.08
    _add_large_visual(slide, icon_key, icon_x, icon_y, icon_w, icon_h)
    text_y = icon_y + icon_h + 0.01
    text_h = max(0.30, y + h - text_y - 0.06)
    _add_text(slide, title, x + 0.08, text_y, w - 0.16, min(0.34, text_h * 0.55), size=title_size, bold=True, color=C["ink"], align=PP_ALIGN.CENTER)
    if subtitle and text_h > 0.42:
        _add_text(slide, subtitle, x + 0.08, text_y + 0.30, w - 0.16, max(0.18, text_h - 0.30), size=subtitle_size, color=C["muted"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
    return card


def _status_icon_key(status: str, test_code: Any = None) -> str:
    value = _safe_text(status, "", 80).upper()
    if test_code is not None:
        return "success" if test_code == 0 else "error"
    if any(x in value for x in ("FAIL", "ERROR", "FAILED")):
        return "error"
    if any(x in value for x in ("WARN", "PENDING", "WAIT", "RUNNING")):
        return "warning"
    if any(x in value for x in ("PASS", "SUCCESS", "COMPLETED", "DONE")):
        return "success"
    return "generic_component"


def _rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _safe_text(value: Any, fallback: str = "", limit: int = 500) -> str:
    if value is None:
        return fallback
    if isinstance(value, (dict, list, tuple)):
        try:
            text = json.dumps(value, ensure_ascii=False, default=str)
        except Exception:
            text = str(value)
    else:
        text = str(value)
    text = text.replace("\\n", " ").replace("\x00", " ")
    text = re.sub(r"\s+", " ", text).strip()
    text = _redact(text)
    if not text:
        return fallback
    if len(text) > limit:
        return text[: max(0, limit - 1)].rstrip() + "…"
    return text


def _redact(text: str) -> str:
    value = str(text or "")
    value = re.sub(r"\bsk-[A-Za-z0-9_-]{16,}\b", "[REDACTED_TOKEN]", value)
    value = re.sub(r"\bAIza[0-9A-Za-z_-]{20,}\b", "[REDACTED_TOKEN]", value)
    value = re.sub(r"\bgh[pousr]_[A-Za-z0-9]{20,}\b", "[REDACTED_TOKEN]", value)
    value = re.sub(r"\bxox[baprs]-[A-Za-z0-9-]{16,}\b", "[REDACTED_TOKEN]", value)
    value = re.sub(
        r"((?:postgres(?:ql)?|mysql|mariadb|redis):\/\/[^:\s\/@]+:)([^@\s\/]+)(@)",
        r"\1[REDACTED]\3",
        value,
        flags=re.IGNORECASE,
    )
    value = re.sub(
        r"(?im)(\b(?:OPENAI_API_KEY|TAVILY_API_KEY|GEMINI_API_KEY|LANGCHAIN_API_KEY|SUPABASE_DB_PASSWORD|PGPASSWORD|[A-Z0-9_]*(?:API_KEY|TOKEN|SECRET|PASSWORD|PASSWD|PWD))\s*=\s*)([^\r\n;]+)",
        r"\1[REDACTED]",
        value,
    )
    return value


def _items(value: Any) -> list[Any]:
    if value is None or value == "":
        return []
    if isinstance(value, list):
        return [x for x in value if x is not None]
    if isinstance(value, tuple):
        return [x for x in value if x is not None]
    return [value]


def _label(item: Any, fallback: str = "") -> str:
    if isinstance(item, str):
        return _safe_text(item, fallback, 100)
    if isinstance(item, dict):
        for key in ("label", "name", "component", "title", "path", "capability", "id"):
            if item.get(key):
                return _safe_text(item.get(key), fallback, 100)
    return _safe_text(item, fallback, 100)


def _detail(item: Any, fallback: str = "") -> str:
    if isinstance(item, dict):
        for key in ("description", "purpose", "reason", "type", "status"):
            if item.get(key):
                return _safe_text(item.get(key), fallback, 180)
    return fallback


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    bold: bool = False,
    color: str = C["ink"],
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    font: str = FONT,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _safe_text(text, "", 5000)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = C["white"],
    line: str = C["line"],
    radius: bool = True,
    line_width: float = 1.0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def _set_shape_text(
    shape,
    title: str,
    subtitle: str = "",
    *,
    title_size: float = 15,
    subtitle_size: float = 9.5,
    title_color: str = C["ink"],
    subtitle_color: str = C["muted"],
    align: PP_ALIGN = PP_ALIGN.CENTER,
):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)

    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _safe_text(title, "", 180)
    run.font.name = FONT
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = _rgb(title_color)

    if subtitle:
        p2 = frame.add_paragraph()
        p2.alignment = align
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = _safe_text(subtitle, "", 260)
        run2.font.name = FONT
        run2.font.size = Pt(subtitle_size)
        run2.font.color.rgb = _rgb(subtitle_color)


def _add_arrow(slide, x: float, y: float, w: float, h: float, direction: str = "right", color: str = C["blue"]):
    shape_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color)
    shp.line.fill.background()
    return shp


def _add_title(slide, title: str, subtitle: str = "", eyebrow: str = "THEANOVA AGENTSTUDIO"):
    _add_text(slide, eyebrow, 0.55, 0.25, 8.7, 0.28, size=9.5, bold=True, color=C["blue"])
    _add_text(slide, title, 0.55, 0.53, 11.8, 0.52, size=25, bold=True, color=C["ink"])
    if subtitle:
        _add_text(slide, subtitle, 0.57, 1.04, 11.9, 0.42, size=10.5, color=C["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.48), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(C["line"])
    line.line.fill.background()


def _add_footer(slide, project_name: str, version: str, page: int):
    _add_text(slide, project_name, 0.55, 7.18, 7.8, 0.18, size=7.5, color=C["muted"])
    _add_text(slide, f"THEANOVA AgentStudio v{version}  ·  {page}", 9.25, 7.18, 3.5, 0.18, size=7.5, color=C["muted"], align=PP_ALIGN.RIGHT)


def _metric(slide, x: float, y: float, w: float, label: str, value: str, sub: str = "", tone: str = "blue"):
    tone_map = {
        "blue": (C["blue_soft"], C["blue"]),
        "green": (C["green_soft"], C["green"]),
        "orange": (C["orange_soft"], C["orange"]),
        "red": (C["red_soft"], C["red"]),
        "purple": (C["purple_soft"], C["purple"]),
        "cyan": (C["cyan_soft"], C["cyan"]),
    }
    fill, accent = tone_map.get(tone, tone_map["blue"])
    card = _add_box(slide, x, y, w, 1.0, fill=fill, line=fill)
    # small accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(accent); bar.line.fill.background()
    _add_text(slide, label, x + 0.18, y + 0.09, w - 0.28, 0.22, size=8.5, bold=True, color=C["muted"])
    _add_text(slide, value, x + 0.18, y + 0.30, w - 0.28, 0.34, size=17, bold=True, color=C["ink"])
    if sub:
        _add_text(slide, sub, x + 0.18, y + 0.68, w - 0.28, 0.20, size=7.8, color=C["muted"])
    return card


def _bullet_list(slide, title: str, items: Iterable[str], x: float, y: float, w: float, h: float, *, accent: str = C["blue"], max_items: int = 7):
    card = _add_box(slide, x, y, w, h, fill=C["white"], line=C["line"])
    _add_text(slide, title, x + 0.18, y + 0.12, w - 0.35, 0.28, size=11, bold=True, color=C["ink"])
    rows = [_safe_text(v, "", 160) for v in list(items) if _safe_text(v, "", 160)]
    rows = rows[:max_items]
    if not rows:
        rows = ["정보가 없습니다."]
    top = y + 0.52
    row_h = max(0.30, min(0.48, (h - 0.65) / max(1, len(rows))))
    for idx, item in enumerate(rows):
        yy = top + idx * row_h
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.20), Inches(yy + 0.095), Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(accent); dot.line.fill.background()
        _add_text(slide, item, x + 0.34, yy, w - 0.52, row_h, size=8.8, color=C["ink"], valign=MSO_ANCHOR.TOP)
    return card


def _new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = _rgb(C["white"])
    return slide


def _factory_phases(definition: dict[str, Any] | None) -> list[dict[str, Any]]:
    fallback = [
        {"title": "요구 이해", "subtitle": "목표·입력·출력·제약 구조화", "nodes": ["요구사항 분석", "프로젝트 분석"]},
        {"title": "Agent 설계", "subtitle": "기능·도구·구조·업무 흐름 결정", "nodes": ["기능 설계", "Tool / MCP 판단", "Agent 아키텍처", "대상 Workflow", "파일 계획"]},
        {"title": "제작", "subtitle": "코드와 실행 환경 구성", "nodes": ["체크포인트", "실행 승인", "코드 생성 / 수정", "환경 구성"]},
        {"title": "검증 & 완성", "subtitle": "실행·복구·완료 확인", "nodes": ["테스트", "디버그 / 복구", "완성 패키지", "최종 검토"]},
    ]
    raw = (definition or {}).get("factory_phases") or fallback
    result: list[dict[str, Any]] = []
    for phase in _items(raw)[:5]:
        if isinstance(phase, dict):
            nodes = []
            for node in _items(phase.get("nodes")):
                nodes.append(_label(node, "단계"))
            result.append({
                "title": _safe_text(phase.get("title") or phase.get("label") or phase.get("id"), "단계", 80),
                "subtitle": _safe_text(phase.get("subtitle") or phase.get("description"), "", 120),
                "nodes": nodes[:6],
            })
        else:
            result.append({"title": _safe_text(phase, "단계", 80), "subtitle": "", "nodes": []})
    return result or fallback


def _add_cover(prs: Presentation, payload: dict[str, Any], version: str):
    slide = _new_slide(prs)
    project = _safe_text(payload.get("project_name"), "AgentStudio Project", 120)
    generated_at = _safe_text(payload.get("generated_at"), datetime.now().isoformat(timespec="seconds"), 80)
    scope = str(payload.get("scope") or "ALL").upper()

    # Decorative, fully editable native shapes.
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(7.5))
    band.fill.solid(); band.fill.fore_color.rgb = _rgb("F8FBFF"); band.line.fill.background()
    # remove/re-add title shapes above background; background inserted late here, so send to back isn't exposed.
    # Instead use semi-transparent-feel solids on top with large white center card.
    hero = _add_box(slide, 0.65, 0.62, 12.0, 6.15, fill=C["white"], line=C["line"])
    for i, (x, y, w, h, color) in enumerate([
        (8.95, 0.92, 2.9, 1.15, C["blue_soft"]),
        (9.55, 2.25, 2.45, 1.15, C["purple_soft"]),
        (8.75, 3.72, 3.25, 1.15, C["cyan_soft"]),
        (9.55, 5.08, 2.45, 0.85, C["green_soft"]),
    ]):
        shp = _add_box(slide, x, y, w, h, fill=color, line=color)
        _set_shape_text(shp, ["Workflow", "Execution", "Architecture", "Report"][i], "editable PowerPoint objects", title_size=13, subtitle_size=7.5)

    _add_text(slide, "THEANOVA AGENTSTUDIO", 1.10, 1.12, 6.5, 0.30, size=10, bold=True, color=C["blue"])
    _add_text(slide, project, 1.10, 1.62, 7.3, 0.92, size=30, bold=True, color=C["ink"], valign=MSO_ANCHOR.TOP)
    report = payload.get("report") or {}
    profile = report.get("projectProfile") or {}
    project_type = _safe_text(profile.get("project_type_label"), "", 90)
    cover_subtitle = f"{project_type} · Workflow · 실행 · 분석 · 아키텍처 리포트" if project_type else "Agent 개발 · 실행 · 분석 · 아키텍처 리포트"
    _add_text(slide, cover_subtitle, 1.12, 2.62, 7.2, 0.45, size=15, color=C["navy2"])
    _add_text(slide, f"내보내기 범위  {_SCOPE_LABEL.get(scope, scope)}", 1.12, 3.45, 4.0, 0.32, size=10, bold=True, color=C["purple"])
    _add_text(slide, f"생성 시각  {generated_at}", 1.12, 3.88, 5.6, 0.28, size=9.5, color=C["muted"])
    _add_text(slide, "모든 핵심 박스·텍스트·연결 요소는 PowerPoint에서 직접 수정할 수 있도록 네이티브 도형으로 생성됩니다.", 1.12, 4.52, 6.8, 0.92, size=11, color=C["ink"], valign=MSO_ANCHOR.TOP)
    _add_text(slide, f"THEANOVA AgentStudio v{version}", 1.12, 5.82, 6.4, 0.30, size=9, bold=True, color=C["muted"])
    return slide


def _add_factory_workflow_slide(prs: Presentation, payload: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    _add_title(slide, "AgentStudio 제작 Workflow", "자연어 요구를 실행 가능한 Agent 프로그램으로 만드는 전체 제작 공정", "WORKFLOW · AGENT FACTORY")
    _add_large_visual(slide, "workflow", 11.62, 0.34, 0.92, 0.92)
    phases = _factory_phases(payload.get("workflow_definition") or {})
    count = len(phases)
    left = 0.65
    gap = 0.16
    total_w = 12.03
    card_w = (total_w - gap * (count - 1)) / max(1, count)
    y = 2.10
    phase_colors = [(C["blue_soft"], C["blue"]), (C["purple_soft"], C["purple"]), (C["cyan_soft"], C["cyan"]), (C["green_soft"], C["green"]), (C["orange_soft"], C["orange"])]
    for idx, phase in enumerate(phases):
        x = left + idx * (card_w + gap)
        fill, accent = phase_colors[idx % len(phase_colors)]
        card = _add_box(slide, x, y, card_w, 3.25, fill=fill, line=accent)
        _add_text(slide, f"{idx + 1:02d}", x + 0.16, y + 0.13, 0.48, 0.26, size=9, bold=True, color=accent)
        _add_text(slide, phase["title"], x + 0.16, y + 0.42, card_w - 0.30, 0.40, size=14, bold=True, color=C["ink"])
        if phase.get("subtitle"):
            _add_text(slide, phase["subtitle"], x + 0.16, y + 0.80, card_w - 0.30, 0.50, size=8.5, color=C["muted"], valign=MSO_ANCHOR.TOP)
        node_y = y + 1.44
        for node_idx, node in enumerate(phase.get("nodes") or []):
            if node_idx >= 5:
                break
            node_box = _add_box(slide, x + 0.16, node_y + node_idx * 0.33, card_w - 0.32, 0.27, fill=C["white"], line=C["line"])
            _set_shape_text(node_box, node, "", title_size=8.3, title_color=C["ink"], align=PP_ALIGN.LEFT)
        if idx < count - 1:
            _add_arrow(slide, x + card_w + 0.025, y + 1.45, gap - 0.05, 0.28, "right", C["navy2"])

    repair = _add_box(slide, 0.95, 5.62, 11.45, 0.82, fill=C["orange_soft"], line=C["orange"])
    _set_shape_text(repair, "자동 복구 루프", "TEST → DEBUG → CODE → ENV → RE-TEST  ·  실패 원인을 분석해 다시 수정한 뒤 재검증", title_size=11, subtitle_size=8.2, title_color=C["orange"])
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide


def _normalize_steps(workflow: dict[str, Any] | None) -> list[dict[str, str]]:
    rows: list[dict[str, str]] = []
    for idx, step in enumerate(_items((workflow or {}).get("steps"))):
        if isinstance(step, str):
            rows.append({"label": _safe_text(step, f"Step {idx + 1}", 90), "description": "", "type": "process"})
        elif isinstance(step, dict):
            rows.append({
                "label": _safe_text(step.get("label") or step.get("name"), f"Step {idx + 1}", 90),
                "description": _safe_text(step.get("description"), "", 160),
                "type": _safe_text(step.get("type"), "process", 50),
            })
    return rows


def _add_target_workflow_slide(prs: Presentation, payload: dict[str, Any], report: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    workflow = report.get("targetWorkflow") or report.get("target_workflow") or {}
    wf_name = _safe_text(workflow.get("name"), "개발 대상 Agent Workflow", 100)
    _add_title(slide, wf_name, "사용자가 Agent를 실행했을 때 처리되는 실제 업무 흐름", "WORKFLOW · TARGET AGENT")
    _add_large_visual(slide, "agent_cube", 11.62, 0.34, 0.92, 0.92)
    steps = _normalize_steps(workflow)
    if not steps:
        _bullet_list(slide, "Workflow", ["아직 대상 Agent Workflow가 생성되지 않았습니다."], 1.0, 2.2, 11.3, 2.2, accent=C["purple"])
    else:
        max_steps = 8
        show = steps[:max_steps]
        n = len(show)
        start_x = 0.65
        gap = 0.12
        total_w = 12.05
        box_w = min(1.55, (total_w - gap * (n - 1)) / n)
        if n <= 5:
            box_w = 2.0
            total_used = box_w * n + gap * (n - 1)
            start_x = (13.333 - total_used) / 2
        y = 2.42
        for idx, row in enumerate(show):
            x = start_x + idx * (box_w + gap)
            tone_fill, tone_line = (C["purple_soft"], C["purple"]) if idx not in (0, n - 1) else ((C["blue_soft"], C["blue"]) if idx == 0 else (C["green_soft"], C["green"]))
            card = _add_box(slide, x, y, box_w, 1.85, fill=tone_fill, line=tone_line)
            _add_text(slide, f"{idx + 1:02d}", x + 0.12, y + 0.11, 0.40, 0.22, size=8.5, bold=True, color=tone_line)
            _add_text(slide, row["label"], x + 0.12, y + 0.39, box_w - 0.24, 0.52, size=11.5, bold=True, color=C["ink"], align=PP_ALIGN.CENTER)
            if row["description"]:
                _add_text(slide, row["description"], x + 0.12, y + 0.96, box_w - 0.24, 0.70, size=7.8, color=C["muted"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
            if idx < n - 1:
                _add_arrow(slide, x + box_w + 0.01, y + 0.79, max(0.08, gap - 0.02), 0.24, "right", C["navy2"])
        if len(steps) > max_steps:
            _add_text(slide, f"※ 전체 {len(steps)}개 단계 중 핵심 {max_steps}개 단계를 표시했습니다.", 0.8, 4.52, 6.8, 0.25, size=8, color=C["muted"])

    branches = _items(workflow.get("branches"))
    retry = workflow.get("retry") or workflow.get("retry_policy") or workflow.get("failure_policy") or ""
    _metric(slide, 0.82, 5.25, 2.6, "Workflow 단계", f"{len(steps)}개", "실제 처리 단계", "blue")
    _metric(slide, 3.62, 5.25, 2.6, "분기", f"{len(branches)}개", "조건 분기", "purple")
    _metric(slide, 6.42, 5.25, 2.6, "재시도 / 실패 처리", "적용" if retry else "정보 없음", _safe_text(retry, "정책 정보 없음", 60), "orange")
    db_plan = report.get("databasePlan") or {}
    _metric(slide, 9.22, 5.25, 2.6, "DB 설계", "사용" if db_plan.get("enabled") else "미사용", f"Table {len(_items(db_plan.get('tables')))}개", "cyan")
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide


def _add_execution_slide(prs: Presentation, payload: dict[str, Any], report: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    profile = report.get("projectProfile") or {}
    adaptive = bool(profile.get("project_type")) and not _items(report.get("createdFiles")) and not _items(report.get("modifiedFiles"))
    exec_subtitle = (
        f"{_safe_text(profile.get('project_type_label'), '현재 프로젝트', 80)} · 실행 준비 상태와 실제 테스트 결과"
        if adaptive else
        "Agent 제작 Workflow의 실행·테스트·파일 변경·디버그 상태 요약"
    )
    _add_title(slide, "실행 결과", exec_subtitle, "PROJECT · EXECUTION" if adaptive else "AGENT FACTORY · EXECUTION")
    _add_large_visual(slide, "terminal", 11.62, 0.34, 0.92, 0.92)
    status = _safe_text(report.get("status"), "NOT_STARTED", 80)
    test_code = report.get("testReturncode")
    test_known = test_code is not None
    test_ok = test_code == 0
    created = _items(report.get("createdFiles"))
    modified = _items(report.get("modifiedFiles"))
    debug = _items(report.get("debugHistory"))

    _metric(slide, 0.65, 1.80, 2.85, "개발 상태", status, "현재 Workflow 상태", "green" if "COMPLETED" in status else "blue")
    _metric(slide, 3.72, 1.80, 2.85, "테스트", "PASS" if test_ok else ("FAIL" if test_known else "대기"), f"Exit Code {test_code if test_known else '-'}", "green" if test_ok else ("red" if test_known else "orange"))
    _metric(slide, 6.79, 1.80, 2.85, "파일 변경", f"생성 {len(created)}", f"수정 {len(modified)}", "cyan")
    _metric(slide, 9.86, 1.80, 2.85, "디버그 / 복구", f"{int(report.get('debugIteration') or 0)}회", "자동 복구 수행" if debug else "재시도 없음", "orange")

    command = _safe_text(report.get("testCommand"), "python -m compileall .", 220)
    output = _safe_text((report.get("testResult") or {}).get("output"), "테스트 출력이 없습니다.", 950)
    created_labels = [_label(x, _safe_text(x, "", 100)) for x in created]
    modified_labels = [_label(x, _safe_text(x, "", 100)) for x in modified]

    left = _add_box(slide, 0.65, 3.10, 6.1, 3.45, fill=C["panel"], line=C["line"])
    _add_text(slide, "테스트 실행", 0.87, 3.26, 2.5, 0.28, size=11.5, bold=True, color=C["ink"])
    _add_text(slide, command, 0.87, 3.66, 5.65, 0.38, size=8.6, bold=True, color=C["navy2"], font=FONT_MONO)
    log_box = _add_box(slide, 0.87, 4.15, 5.65, 2.12, fill=C["white"], line=C["line"])
    _add_text(slide, output, 1.02, 4.29, 5.35, 1.84, size=7.4, color=C["ink"], font=FONT_MONO, valign=MSO_ANCHOR.TOP)

    _add_box(slide, 6.96, 3.10, 5.75, 3.45, fill=C["white"], line=C["line"])
    _add_text(slide, "생성 / 수정 파일", 7.18, 3.26, 2.8, 0.28, size=11.5, bold=True)
    rows = [f"+ {_safe_text(x, '', 92)}" for x in created_labels[:6]] + [f"~ {_safe_text(x, '', 92)}" for x in modified_labels[:6]]
    if not rows:
        rows = ["파일 변경 정보가 없습니다."]
    _bullet_list(slide, "변경 목록", rows, 7.15, 3.70, 5.30, 2.57, accent=C["cyan"], max_items=9)
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide


def _add_report_slide(prs: Presentation, payload: dict[str, Any], report: dict[str, Any], style: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    profile = report.get("projectProfile") or {}
    adaptive = bool(profile.get("project_type"))
    report_subtitle = (
        f"{_safe_text(profile.get('project_type_label'), '현재 프로젝트', 80)} · Workflow · Architecture · Tool · Data 구성"
        if adaptive else
        "요구사항 · Architecture · MCP / Tool · Workflow · 코드 품질 · 최종 상태"
    )
    _add_title(slide, "분석 리포트", report_subtitle, "PROJECT · ADAPTIVE REPORT" if adaptive else "AGENT DEVELOPMENT · REPORT")
    _add_large_visual(slide, "report", 11.62, 0.34, 0.92, 0.92)
    workflow = report.get("targetWorkflow") or {}
    mcp = _items((report.get("toolMcpPlan") or {}).get("decisions"))
    caps = _items((report.get("capabilityPlan") or {}).get("capabilities"))
    status = _safe_text(report.get("status"), "NOT_STARTED", 80)

    _metric(slide, 0.65, 1.75, 2.85, "Workflow 단계", f"{len(_normalize_steps(workflow))}개", f"분기 {len(_items(workflow.get('branches')))}개", "blue")
    _metric(slide, 3.72, 1.75, 2.85, "MCP / Tool", f"{len(mcp)}개", "연결 판단 결과", "purple")
    _metric(slide, 6.79, 1.75, 2.85, "코딩 스타일", "PASS" if int(style.get("fail") or 0) == 0 else "FAIL", f"경고 {int(style.get('warning') or 0)} · 오류 {int(style.get('fail') or 0)}", "green" if int(style.get("fail") or 0) == 0 else "red")
    _metric(slide, 9.86, 1.75, 2.85, "최종 상태", status, "Agent Factory", "green" if "COMPLETED" in status else "orange")

    req = report.get("requirementSpec") or {}
    goal = _safe_text(req.get("goal") or payload.get("workflow_request"), "요구사항 정보가 없습니다.", 520)
    _bullet_list(slide, "요구사항 / 목표", [goal, f"Acceptance Criteria {len(_items(req.get('acceptance_criteria')))}개", f"제약 조건 {len(_items(req.get('constraints')))}개"], 0.65, 3.05, 5.85, 2.95, accent=C["blue"], max_items=4)

    cap_labels = [_label(x, f"Capability {i + 1}") for i, x in enumerate(caps)]
    _bullet_list(slide, "Capabilities", cap_labels, 6.72, 3.05, 2.90, 2.95, accent=C["purple"], max_items=6)

    mcp_labels = []
    for idx, item in enumerate(mcp):
        if isinstance(item, dict):
            mcp_labels.append(f"{_safe_text(item.get('capability'), f'Capability {idx+1}', 70)} · {_safe_text(item.get('execution_type'), 'none', 36)}")
        else:
            mcp_labels.append(_safe_text(item, "", 100))
    _bullet_list(slide, "MCP / Tool 결정", mcp_labels, 9.85, 3.05, 2.86, 2.95, accent=C["cyan"], max_items=6)

    arch = report.get("architecture") or {}
    _add_text(slide, f"Architecture  Components {len(_items(arch.get('components')))} · Interfaces {len(_items(arch.get('interfaces')))} · Persistence {len(_items(arch.get('persistence')))} · Security {len(_items(arch.get('security')))}", 0.68, 6.24, 11.9, 0.26, size=8.3, bold=True, color=C["navy2"])
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide


def _add_layer_label(slide, text: str, x: float, y: float, w: float, h: float, fill: str, color: str):
    shp = _add_box(slide, x, y, w, h, fill=fill, line=fill)
    _set_shape_text(shp, text, "", title_size=9.5, title_color=color, align=PP_ALIGN.CENTER)
    return shp


def _add_architecture_slide(prs: Presentation, payload: dict[str, Any], report: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    req = report.get("requirementSpec") or {}
    profile = report.get("projectProfile") or {}
    arch = report.get("architecture") or {}
    source = _safe_text(arch.get("source"), "", 60)
    adaptive = source == "PROJECT_SOURCE_INFERENCE"
    title = "프로젝트 적응형 아키텍처" if adaptive else "Target Project Architecture"
    goal = _safe_text(
        req.get("goal") or profile.get("summary"),
        "현재 프로젝트에서 확인된 구성 요소를 기반으로 Architecture를 시각화합니다.",
        220,
    )
    type_label = _safe_text(profile.get("project_type_label"), "", 90)
    subtitle = goal if not type_label else f"{type_label} · {goal}"
    _add_title(slide, title, subtitle, "ARCHITECTURE · PROJECT ADAPTIVE · LARGE VISUAL")

    components = _items(arch.get("components"))
    interfaces = _items(arch.get("interfaces"))
    persistence = _items(arch.get("persistence"))
    security = _items(arch.get("security"))
    states = _items(arch.get("state"))
    infrastructure = _items(arch.get("infrastructure"))

    # Interface layer: only items actually present in the project snapshot.
    _add_layer_label(slide, "INTERFACE / ENTRY", 0.55, 1.70, 1.30, 1.08, C["blue_soft"], C["blue"])
    interface_rows = interfaces[:4]
    if not interface_rows:
        interface_rows = [{"label": "Program Entry Point", "description": "프로젝트 실행 진입점"}]
    n = len(interface_rows)
    card_w = min(2.35, (10.75 - 0.20 * max(0, n - 1)) / max(1, n))
    total = card_w * n + 0.20 * max(0, n - 1)
    start_x = 1.98 + max(0, (10.35 - total) / 2)
    for idx, item in enumerate(interface_rows):
        label = _label(item, f"Interface {idx+1}")
        detail = _detail(item, "Project interface")
        icon = _icon_key_for_text(label, detail, fallback="api_gateway")
        x = start_x + idx * (card_w + 0.20)
        _add_visual_card(slide, x, 1.66, card_w, 1.16, title=label, subtitle=detail, icon_key=icon, fill=C["blue_soft"], line=C["blue"], icon_ratio=0.58, title_size=8.8, subtitle_size=6.1)
        if idx < n - 1:
            _add_arrow(slide, x + card_w + 0.02, 2.11, 0.15, 0.16, "right", C["navy2"])

    # Core layer: real detected/generated components only; no fixed AgentStudio stack injected.
    _add_layer_label(slide, "CORE / SERVICE", 0.55, 3.02, 1.30, 1.72, C["purple_soft"], C["purple"])
    core_rows = components[:5]
    if not core_rows:
        core_rows = [{"label": "Application Core", "description": "프로젝트 핵심 로직"}]
    n = len(core_rows)
    card_w = min(1.95, (10.78 - 0.18 * max(0, n - 1)) / max(1, n))
    total = card_w * n + 0.18 * max(0, n - 1)
    start_x = 1.98 + max(0, (10.35 - total) / 2)
    for idx, item in enumerate(core_rows):
        label = _label(item, f"Component {idx+1}")
        detail = _detail(item, "Project component")
        icon = _icon_key_for_text(label, detail, fallback="generic_component")
        x = start_x + idx * (card_w + 0.18)
        _add_visual_card(slide, x, 2.98, card_w, 1.72, title=label, subtitle=detail, icon_key=icon, fill=C["purple_soft"], line=C["purple"], icon_ratio=0.56, title_size=8.5, subtitle_size=6.0)
        if idx < n - 1:
            _add_arrow(slide, x + card_w + 0.015, 3.69, 0.15, 0.16, "right", C["navy2"])

    state_text = f"State {len(states)} · Security {len(security)}"
    state_band = _add_box(slide, 2.25, 4.79, 9.45, 0.38, fill=C["panel"], line=C["line"])
    _set_shape_text(state_band, "State / Policy", state_text, title_size=8.1, subtitle_size=6.1)

    # Data/runtime layer: use only actual persistence/infrastructure detections.
    _add_layer_label(slide, "DATA / RUNTIME", 0.55, 5.31, 1.30, 1.21, C["green_soft"], C["green"])
    data_rows = (persistence + infrastructure)[:5]
    if data_rows:
        n = len(data_rows)
        card_w = min(1.95, (10.78 - 0.18 * max(0, n - 1)) / max(1, n))
        total = card_w * n + 0.18 * max(0, n - 1)
        start_x = 1.98 + max(0, (10.35 - total) / 2)
        for idx, item in enumerate(data_rows):
            label = _label(item, f"Runtime {idx+1}")
            detail = _detail(item, "Persistence / Runtime")
            icon = _icon_key_for_text(label, detail, fallback="storage")
            x = start_x + idx * (card_w + 0.18)
            _add_visual_card(slide, x, 5.27, card_w, 1.23, title=label, subtitle=detail, icon_key=icon, fill=C["green_soft"], line=C["green"], icon_ratio=0.57, title_size=8.1, subtitle_size=5.9)
            _add_arrow(slide, x + card_w/2 - 0.10, 5.00, 0.20, 0.25, "down", C["green"])
    else:
        empty = _add_box(slide, 2.18, 5.38, 9.62, 0.82, fill=C["panel"], line=C["line"])
        _set_shape_text(empty, "별도 Persistence / Infrastructure 미감지", "현재 소스에서 확인되지 않은 DB·Cache·Cloud 기술은 임의로 표시하지 않습니다.", title_size=10.2, subtitle_size=7.0, title_color=C["muted"])

    _add_text(slide, "Project Adaptive: 실제 소스/Agent 설계에서 확인된 요소만 사용하며, 감지되지 않은 기술은 PPT에 자동 삽입하지 않습니다.", 1.55, 6.67, 10.7, 0.22, size=7.2, color=C["muted"], align=PP_ALIGN.CENTER)
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide

def _add_asbuilt_slide(prs: Presentation, payload: dict[str, Any], report: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    _add_title(slide, "As-Built Architecture & Conformance", "실제 생성 코드에서 감지한 구조와 Design Architecture 일치 여부", "ARCHITECTURE · AS-BUILT")
    arch = report.get("asBuiltArchitecture") or {}
    conf = report.get("architectureConformance") or {}
    components = _items(arch.get("components"))
    interfaces = _items(arch.get("interfaces"))
    persistence = _items(arch.get("persistence"))
    scan = arch.get("scan") or {}
    score = float(conf.get("score") or 0)
    ok = conf.get("ok") is True

    _metric(slide, 0.65, 1.75, 2.85, "소스 파일", f"{int(scan.get('source_file_count') or 0)}개", "정적 분석 대상", "blue")
    _metric(slide, 3.72, 1.75, 2.85, "감지 Component", f"{len(components)}개", "실제 구현 구조", "purple")
    _metric(slide, 6.79, 1.75, 2.85, "Interface / Persistence", f"{len(interfaces)} / {len(persistence)}", "실제 소스 증거", "cyan")
    _metric(slide, 9.86, 1.75, 2.85, "Conformance", f"{score:.0f}점" if conf else "대기", "PASS" if ok else ("검증 필요" if conf else "분석 전"), "green" if ok else "orange")

    comp_rows = []
    for idx, item in enumerate(components[:8]):
        status = _safe_text(item.get("status"), "detected", 30) if isinstance(item, dict) else "detected"
        comp_rows.append(f"{_label(item, f'Component {idx+1}')} · {status}")
    _bullet_list(slide, "실제 감지 구성 요소", comp_rows, 0.65, 3.05, 5.85, 3.15, accent=C["purple"], max_items=8)

    mismatches = _items(conf.get("mismatches"))
    mismatch_rows = []
    for idx, item in enumerate(mismatches[:8]):
        if isinstance(item, dict):
            mismatch_rows.append(_safe_text(item.get("message") or item.get("expected") or item.get("component") or item, f"Mismatch {idx+1}", 145))
        else:
            mismatch_rows.append(_safe_text(item, f"Mismatch {idx+1}", 145))
    if not mismatch_rows:
        mismatch_rows = ["Critical 누락 없음" if ok else "비교 결과가 아직 생성되지 않았습니다."]
    _bullet_list(slide, "Design ↔ As-Built 비교", mismatch_rows, 6.72, 3.05, 5.99, 3.15, accent=C["green"] if ok else C["orange"], max_items=8)
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide


def _add_project_stack_slide(prs: Presentation, payload: dict[str, Any], report: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    profile = report.get("projectProfile") or {}
    arch = report.get("architecture") or {}
    project_type = _safe_text(profile.get("project_type_label"), "Project Runtime", 90)
    tech_stack = _items(profile.get("tech_stack"))
    if not tech_stack:
        tech_stack = _items(profile.get("languages"))
    infrastructure = _items(arch.get("infrastructure")) or _items(profile.get("infrastructure"))
    persistence = _items(arch.get("persistence"))

    _add_title(slide, "Project Technology & Runtime", f"{project_type} · 실제 프로젝트에서 감지된 기술/데이터/인프라", "ARCHITECTURE · PROJECT STACK · ADAPTIVE")

    rows = tech_stack[:6]
    if not rows:
        rows = ["Application Source"]
    count = len(rows)
    gap = 0.18
    total_w = 11.85
    card_w = min(1.82, (total_w - gap * max(0, count - 1)) / max(1, count))
    used = card_w * count + gap * max(0, count - 1)
    start_x = (13.333 - used) / 2
    for idx, item in enumerate(rows):
        label = _label(item, f"Technology {idx+1}")
        icon = _icon_key_for_text(label, fallback="code")
        x = start_x + idx * (card_w + gap)
        _add_visual_card(slide, x, 1.72, card_w, 2.02, title=label, subtitle="Detected Technology", icon_key=icon, fill=C["blue_soft"] if idx % 2 == 0 else C["purple_soft"], line=C["blue"] if idx % 2 == 0 else C["purple"], icon_ratio=0.60, title_size=9.0, subtitle_size=6.0)

    _add_text(slide, "DATA / STATE", 0.78, 4.08, 1.4, 0.25, size=9.2, bold=True, color=C["green"])
    data = persistence[:5]
    if data:
        for idx, item in enumerate(data):
            x = 1.06 + idx * 2.20
            label = _label(item, f"Data {idx+1}")
            detail = _detail(item, "Detected persistence")
            _add_visual_card(slide, x, 4.37, 1.84, 1.62, title=label, subtitle=detail, icon_key=_icon_key_for_text(label, detail, fallback="database"), fill=C["green_soft"], line=C["green"], icon_ratio=0.57, title_size=8.3, subtitle_size=5.8)
    else:
        _add_text(slide, "별도 DB / Cache / Vector Store가 소스에서 감지되지 않았습니다.", 1.08, 4.55, 5.2, 0.35, size=9.2, color=C["muted"])

    _add_text(slide, "INFRASTRUCTURE", 7.03, 4.08, 1.7, 0.25, size=9.2, bold=True, color=C["cyan"])
    infra = infrastructure[:4]
    if infra:
        for idx, item in enumerate(infra):
            x = 7.08 + idx * 1.38
            label = _label(item, f"Infra {idx+1}")
            _add_large_visual(slide, _icon_key_for_text(label, _detail(item, ""), fallback="server"), x, 4.48, 0.86, 0.82)
            _add_text(slide, label, x - 0.18, 5.34, 1.22, 0.48, size=7.4, bold=True, color=C["ink"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
    else:
        _add_text(slide, "Container / Cloud / Cluster 인프라가 소스에서 감지되지 않았습니다.", 7.10, 4.55, 4.8, 0.35, size=9.0, color=C["muted"])

    _add_text(slide, "※ 이 슬라이드는 AgentStudio 자체 기술 스택이 아니라 현재 대상 프로젝트에서 실제 감지된 기술만 표시합니다.", 1.15, 6.55, 11.0, 0.26, size=7.4, color=C["muted"], align=PP_ALIGN.CENTER)
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide

def _add_platform_architecture_slide(prs: Presentation, payload: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    _add_title(slide, "THEANOVA AgentStudio 플랫폼 아키텍처", "큰 대표 그래픽 중심의 Client → Interface → Security/API → Agent/Service 구조", "ARCHITECTURE · PLATFORM · LARGE VISUAL")

    # Client devices: intentionally large, visually dominant illustrations.
    client_band = _add_box(slide, 0.65, 1.62, 12.05, 1.12, fill="F8FBFF", line=C["blue"])
    client_items = [
        (1.08, "user_group", "User"),
        (3.65, "laptop", "Desktop / Web"),
        (6.22, "mobile", "Mobile"),
        (8.79, "globe", "External / Internet"),
    ]
    for x, icon, label in client_items:
        _add_large_visual(slide, icon, x, 1.68, 1.02, 0.82)
        _add_text(slide, label, x-0.17, 2.45, 1.36, 0.20, size=7.8, bold=True, color=C["navy2"], align=PP_ALIGN.CENTER)
    _add_text(slide, "CLIENT / USER", 10.65, 1.96, 1.60, 0.28, size=10, bold=True, color=C["blue"], align=PP_ALIGN.CENTER)

    interface = _add_box(slide, 1.45, 2.90, 10.45, 0.72, fill=C["cyan_soft"], line=C["cyan"])
    _add_large_visual(slide, "globe", 2.05, 2.91, 0.76, 0.68)
    _add_text(slide, "Frontend / Interface", 3.04, 3.02, 3.35, 0.27, size=12.5, bold=True, color=C["ink"])
    _add_text(slide, "React + TypeScript + Vite · Workspace / Notebook / SQL / Report", 6.15, 3.04, 4.95, 0.24, size=8.0, color=C["muted"], align=PP_ALIGN.RIGHT)
    _add_arrow(slide, 6.49, 2.68, 0.30, 0.20, "down", C["cyan"])

    security = _add_box(slide, 1.45, 3.76, 10.45, 0.72, fill=C["purple_soft"], line=C["purple"])
    _add_large_visual(slide, "shield", 2.05, 3.77, 0.76, 0.68)
    _add_text(slide, "FastAPI Backend / Security / API", 3.04, 3.88, 4.30, 0.27, size=12.0, bold=True, color=C["ink"])
    _add_text(slide, "API · Jobs · Runtime · Auth · Validation", 7.35, 3.90, 3.75, 0.24, size=8.0, color=C["muted"], align=PP_ALIGN.RIGHT)
    _add_arrow(slide, 6.49, 3.62, 0.30, 0.12, "down", C["purple"])

    service_band = _add_box(slide, 0.65, 4.63, 12.05, 1.90, fill="FBF9FF", line=C["purple"])
    _add_text(slide, "AGENT / SERVICE", 0.87, 4.76, 1.40, 0.25, size=8.8, bold=True, color=C["purple"])
    services = [
        (1.15, "agent_cube", "Agent Orchestrator", "Workflow / Planning"),
        (3.46, "workflow", "LangGraph Workflow", "State / Recovery"),
        (5.77, "llm", "LLM Layer", "OpenAI / Ollama"),
        (8.08, "mcp", "MCP Integration", "Client / Server"),
        (10.39, "tool", "Tools / Runtime", "Local / Python / SQL"),
    ]
    for idx, (x, icon, title, sub) in enumerate(services):
        _add_visual_card(slide, x, 4.92, 1.78, 1.43, title=title, subtitle=sub, icon_key=icon, fill=C["white"], line=C["line"], icon_ratio=0.58, title_size=8.4, subtitle_size=6.0)
        if idx < len(services)-1:
            _add_arrow(slide, x+1.82, 5.53, 0.18, 0.16, "right", C["navy2"])

    _add_text(slide, "Foundation / Data / Infrastructure는 다음 슬라이드에서 큰 시각 자산으로 분리해 가독성을 유지합니다.", 1.30, 6.72, 10.8, 0.22, size=7.4, color=C["muted"], align=PP_ALIGN.CENTER)
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide


def _add_platform_foundation_slide(prs: Presentation, payload: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    _add_title(slide, "AgentStudio Foundation & Infrastructure", "LLM · Execution · Persistence · Project State · Infrastructure를 큰 대표 그래픽으로 분리", "ARCHITECTURE · FOUNDATION · LARGE VISUAL")

    foundation = [
        (0.72, "llm", "LLM Layer", "OpenAI · Ollama · LangChain / LangGraph", C["blue_soft"], C["blue"]),
        (3.88, "tool", "Execution Layer", "MCP · Local Tools · Python / Terminal / SQL", C["purple_soft"], C["purple"]),
        (7.04, "database", "Persistence", "PostgreSQL · Redis · pgvector · SQL DB", C["green_soft"], C["green"]),
        (10.20, "report", "Project State", "Reports · Usage / Cost · Recovery / Checkpoint", C["orange_soft"], C["orange"]),
    ]
    for x, icon, title, sub, fill, line_color in foundation:
        _add_visual_card(slide, x, 1.78, 2.42, 2.36, title=title, subtitle=sub, icon_key=icon, fill=fill, line=line_color, icon_ratio=0.58, title_size=11.0, subtitle_size=7.4)

    infra_band = _add_box(slide, 0.65, 4.45, 12.05, 2.02, fill="F8FBFF", line=C["blue"])
    _add_text(slide, "INFRASTRUCTURE / PLATFORM", 0.88, 4.59, 2.20, 0.25, size=9.0, bold=True, color=C["blue"])
    infra = [
        (1.12, "cloud", "Cloud"),
        (3.33, "server", "Server / Container"),
        (5.54, "kubernetes", "Kubernetes / Cluster"),
        (7.75, "storage", "Object / File Storage"),
        (9.96, "network", "Network / External MCP"),
    ]
    for x, icon, label in infra:
        _add_large_visual(slide, icon, x, 4.87, 1.12, 1.05)
        _add_text(slide, label, x-0.20, 5.94, 1.52, 0.28, size=8.3, bold=True, color=C["ink"], align=PP_ALIGN.CENTER)
    _add_footer(slide, _safe_text(payload.get("project_name"), "Project", 100), version, page)
    return slide




def _add_studio_cover(prs: Presentation, payload: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    _add_box(slide, 0, 0, 13.333, 7.5, fill="F7FAFE", line="F7FAFE", radius=False)
    _add_box(slide, 8.35, 0.72, 4.15, 5.95, fill=C["purple_soft"], line=C["purple"])
    _add_large_visual(slide, "agent_cube", 9.18, 1.46, 2.48, 2.20)
    _add_large_visual(slide, "workflow", 8.72, 4.02, 1.22, 1.04)
    _add_large_visual(slide, "database", 10.13, 4.00, 1.22, 1.04)
    _add_large_visual(slide, "tool", 11.52, 4.02, 0.78, 0.98)
    _add_text(slide, "THEANOVA", 0.92, 1.08, 3.0, 0.34, size=13, bold=True, color=C["blue"])
    _add_text(slide, "AgentStudio", 0.92, 1.56, 6.8, 0.78, size=34, bold=True, color=C["ink"])
    _add_text(slide, "Studio Product / Platform Presentation", 0.95, 2.52, 6.6, 0.42, size=16, color=C["navy2"])
    _add_text(slide, "Agent 제작 · Workflow · 코드 편집 · 실행 · 분석 · 아키텍처 · LLM · MCP 통합 개발 환경", 0.95, 3.25, 6.95, 1.15, size=12, color=C["ink"], valign=MSO_ANCHOR.TOP)
    generated_at = _safe_text(payload.get("generated_at"), datetime.now().isoformat(timespec="seconds"), 80)
    _add_text(slide, f"Generated  {generated_at}", 0.95, 4.78, 5.8, 0.30, size=9.2, color=C["muted"])
    _add_text(slide, f"THEANOVA AgentStudio v{version}", 0.95, 5.26, 5.8, 0.30, size=10, bold=True, color=C["purple"])
    _add_footer(slide, "THEANOVA AgentStudio", version, page)
    return slide


def _add_studio_capabilities_slide(prs: Presentation, payload: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    _add_title(slide, "AgentStudio Workspace & Core Capabilities", "Agent를 설계하고 구현·실행·검증하는 Studio의 핵심 작업 영역", "STUDIO · PRODUCT CAPABILITIES")
    rows = [
        ("agent_cube", "Agent 설계", "대화형 요구사항 수집 / Agent 생성", C["blue_soft"], C["blue"]),
        ("workflow", "Workflow", "Agent Factory / Target Workflow", C["purple_soft"], C["purple"]),
        ("code", "코드 편집", "Project File / Notebook / SQL", C["cyan_soft"], C["cyan"]),
        ("terminal", "실행 결과", "Terminal / Test / Recovery", C["green_soft"], C["green"]),
        ("report", "분석 리포트", "Requirement / Quality / Status", C["orange_soft"], C["orange"]),
        ("network", "아키텍처", "Design / As-Built / Conformance", C["purple_soft"], C["purple"]),
        ("llm", "LLM", "OpenAI / Ollama / Routing", C["blue_soft"], C["blue"]),
        ("globe", "Web Browser", "Embedded / Proxy / CDP", C["cyan_soft"], C["cyan"]),
    ]
    positions=[]
    for r in range(2):
        for c in range(4):
            positions.append((0.72+c*3.12, 1.72+r*2.42))
    for (icon,title,sub,fill,line),(x,y) in zip(rows,positions):
        _add_visual_card(slide, x, y, 2.72, 2.03, title=title, subtitle=sub, icon_key=icon, fill=fill, line=line, icon_ratio=0.60, title_size=11.0, subtitle_size=7.0)
    _add_footer(slide, "THEANOVA AgentStudio", version, page)
    return slide


def _add_studio_runtime_slide(prs: Presentation, payload: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    _add_title(slide, "AgentStudio Execution & Runtime", "Frontend → API → Agent Orchestrator → Tool / LLM / Data Runtime", "STUDIO · EXECUTION · RUNTIME")
    chain = [
        ("laptop", "React Workspace", "UI / Editor / Panels", C["blue_soft"], C["blue"]),
        ("api_gateway", "FastAPI Backend", "API / Jobs / Runtime", C["cyan_soft"], C["cyan"]),
        ("agent_cube", "Agent Orchestrator", "Planning / Workflow", C["purple_soft"], C["purple"]),
        ("tool", "Execution Layer", "MCP / Local Tools / Terminal", C["orange_soft"], C["orange"]),
        ("database", "Persistence", "PostgreSQL / State / Reports", C["green_soft"], C["green"]),
    ]
    x0=0.62; gap=0.17; w=2.34; y=2.10
    for idx,(icon,title,sub,fill,line) in enumerate(chain):
        x=x0+idx*(w+gap)
        _add_visual_card(slide,x,y,w,2.32,title=title,subtitle=sub,icon_key=icon,fill=fill,line=line,icon_ratio=0.60,title_size=10.2,subtitle_size=6.8)
        if idx<len(chain)-1:
            _add_arrow(slide,x+w+0.03,y+1.03,gap-0.03,0.24,"right",C["navy2"])
    _add_box(slide, 0.98, 4.92, 11.38, 1.02, fill=C["panel"], line=C["line"])
    _add_text(slide, "실시간 처리", 1.24, 5.12, 1.35, 0.28, size=10.5, bold=True, color=C["ink"])
    _add_text(slide, "WebSocket / SSE · Multi Terminal · Async Jobs · Recovery / Debug · Project State", 2.74, 5.12, 8.75, 0.28, size=9.2, color=C["muted"])
    _add_footer(slide, "THEANOVA AgentStudio", version, page)
    return slide


def _add_studio_governance_slide(prs: Presentation, payload: dict[str, Any], page: int, version: str):
    slide = _new_slide(prs)
    _add_title(slide, "AgentStudio Analysis & Governance", "요구사항부터 코드 품질·복구·As-Built 검증까지 개발 신뢰성을 관리", "STUDIO · ANALYSIS · GOVERNANCE")
    cards=[
        (0.72,1.80,"report","Requirement Analysis","요구사항 / Acceptance Criteria",C["blue_soft"],C["blue"]),
        (3.82,1.80,"shield","Validation / Quality Gate","형식·업무규칙·신뢰도 검증",C["purple_soft"],C["purple"]),
        (6.92,1.80,"code","Coding Governance","Style / Rule / Static Validation",C["cyan_soft"],C["cyan"]),
        (10.02,1.80,"workflow","Recovery Loop","Test → Debug → Repair → Re-test",C["orange_soft"],C["orange"]),
        (2.28,4.28,"network","As-Built Architecture","실제 코드 구조 역분석",C["green_soft"],C["green"]),
        (7.36,4.28,"shield","Conformance","Design ↔ As-Built 일치 검증",C["purple_soft"],C["purple"]),
    ]
    for x,y,icon,title,sub,fill,line in cards:
        _add_visual_card(slide,x,y,2.58,1.80,title=title,subtitle=sub,icon_key=icon,fill=fill,line=line,icon_ratio=0.57,title_size=9.8,subtitle_size=6.5)
    _add_footer(slide, "THEANOVA AgentStudio", version, page)
    return slide



def _set_erd_badge_text(shape, marker: str, color: str):
    """Keep PK/FK/vector badges horizontal even in narrow ERD cards."""
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = marker
    run.font.name = FONT
    run.font.size = Pt(6.4)
    run.font.bold = True
    run.font.color.rgb = _rgb(color)


def _erd_table_card(slide, table: dict[str, Any], x: float, y: float, w: float, h: float, *, accent: str = C["blue"]):
    """Editable table card used by DB ERD slides.

    v5.363 keeps PK/FK markers horizontal and reserves a stable marker column so
    column names/types never squeeze the badge into a vertical two-line label.
    """
    _add_box(slide, x, y, w, h, fill=C["white"], line=accent)
    header_h = 0.42
    _add_box(slide, x, y, w, header_h, fill=accent, line=accent, radius=False)
    schema = _safe_text(table.get("schema"), "", 40)
    name = _safe_text(table.get("name") or table.get("id"), "table", 80)
    title = f"{schema}.{name}" if schema and not name.startswith(f"{schema}.") else name
    title_size = 7.0 if len(title) > 32 else (8.1 if len(title) > 24 else 9.4)
    _add_text(slide, title, x + 0.08, y + 0.03, w - 0.16, 0.32, size=title_size, bold=True, color=C["white"], margin=0.02)
    columns = [item for item in _items(table.get("columns")) if isinstance(item, dict)]
    available_h = max(0.42, h - header_h - 0.14)
    max_rows = max(2, int(available_h / 0.20))
    if len(columns) > max_rows:
        visible_count = max(1, max_rows - 1)
        visible = columns[:visible_count]
        has_more = True
    else:
        visible = columns
        has_more = False
    row_count = max(1, len(visible) + (1 if has_more else 0))
    row_h = min(0.245, available_h / row_count)
    yy = y + header_h + 0.06

    badge_x = x + 0.08
    badge_w = 0.46
    name_x = x + 0.62
    type_x = x + w * 0.61
    type_w = max(0.50, x + w - 0.10 - type_x)
    name_w = max(0.62, type_x - name_x - 0.08)

    for column in visible:
        pk = bool(column.get("primary_key"))
        fk = bool(column.get("foreign_key"))
        vector = bool(column.get("vector")) or "vector" in _safe_text(column.get("data_type"), "", 60).casefold()
        marker = "PK" if pk else ("FK" if fk else ("V" if vector else ""))
        marker_color = C["purple"] if pk else (C["orange"] if fk else (C["cyan"] if vector else C["muted"]))
        if marker:
            badge_h = max(0.14, row_h - 0.055)
            badge = _add_box(slide, badge_x, yy + (row_h - badge_h) / 2, badge_w, badge_h, fill=C["panel"], line=marker_color)
            _set_erd_badge_text(badge, marker, marker_color)
        _add_text(slide, _safe_text(column.get("name"), "column", 55), name_x, yy, name_w, row_h, size=7.25, bold=pk, color=C["ink"], margin=0.015)
        _add_text(slide, _safe_text(column.get("data_type"), "", 55), type_x, yy, type_w, row_h, size=6.45, color=C["muted"], align=PP_ALIGN.RIGHT, margin=0.015)
        yy += row_h
    if has_more:
        _add_text(slide, f"+ {len(columns)-len(visible)} columns", name_x, yy, w - (name_x - x) - 0.10, row_h, size=6.4, color=C["muted"], margin=0.02)


def _add_relational_erd_page(prs: Presentation, payload: dict[str, Any], database: dict[str, Any], tables: list[dict[str, Any]], relationships: list[dict[str, Any]], page: int, version: str, part: int = 1, total_parts: int = 1):
    slide = _new_slide(prs)
    label = _safe_text(database.get("label"), _safe_text(database.get("engine"), "Database", 80), 100)
    suffix = f" · {part}/{total_parts}" if total_parts > 1 else ""
    _add_title(slide, f"DB ERD · {label}{suffix}", "현재 Agent/프로젝트에서 설계 또는 소스로 감지한 DB Schema와 관계", "AGENT · DATABASE ERD" if str((payload.get('db_erd') or {}).get('scope') or 'AGENT').upper() != 'STUDIO' else "STUDIO · DATABASE ERD")
    icon = "vector_db" if str(database.get("kind")) == "vector" else "database"
    _add_large_visual(slide, icon, 11.58, 0.26, 0.78, 0.82)
    _metric(slide, 0.72, 1.67, 2.55, "Tables", str(len(tables)), "현재 슬라이드", "blue")
    _metric(slide, 3.48, 1.67, 2.55, "Relations", str(len([r for r in relationships if isinstance(r, dict)])), "전체 감지 관계", "purple")
    source_raw = _safe_text(database.get("source"), "INFERENCE", 60)
    source_label = {
        "AGENTSTUDIO_SCHEMA_SQL": "STUDIO SCHEMA",
        "AGENT_DATABASE_PLAN": "AGENT PLAN",
        "AGENT_REDIS_PLAN": "AGENT PLAN",
        "PROJECT_SQL_SOURCE": "PROJECT DDL",
        "PROJECT_SOURCE_INFERENCE": "SOURCE SCAN",
        "VECTOR_SCHEMA_INFERENCE": "VECTOR SCAN",
        "TECHNOLOGY_DETECTED": "TECH DETECT",
    }.get(source_raw.upper(), source_raw[:24])
    _metric(slide, 6.24, 1.67, 2.55, "Source", source_label, "Schema 근거", "cyan")
    _metric(slide, 9.00, 1.67, 2.55, "Engine", _safe_text(database.get("engine"), "SQL", 28).upper(), "Database", "green")

    # v5.363: slightly shorter cards create a dedicated routing corridor between ERD rows.
    table_h = 1.70
    table_w = 3.00
    ids = [_safe_text(t.get("id"), "", 100) for t in tables]
    id_set = set(ids)
    parent_score: dict[str, int] = {table_id: 0 for table_id in ids}
    for relation in relationships:
        if not isinstance(relation, dict):
            continue
        target = _safe_text(relation.get("to_table"), "", 100)
        source = _safe_text(relation.get("from_table"), "", 100)
        if target in id_set and source in id_set and target != source:
            parent_score[target] = parent_score.get(target, 0) + 1
    roots = [table_id for table_id in sorted(ids, key=lambda v: (-parent_score.get(v, 0), ids.index(v))) if parent_score.get(table_id, 0) > 0]
    roots = roots[:2]
    root_set = set(roots)
    children = [table_id for table_id in ids if table_id not in root_set]

    positions_by_id: dict[str, tuple[float, float]] = {}
    if roots and children:
        root_xs = [5.16] if len(roots) == 1 else [3.26, 7.06]
        for table_id, x in zip(roots, root_xs):
            positions_by_id[table_id] = (x, 2.82)
        child_count = len(children)
        if child_count == 1:
            child_xs = [5.16]
        elif child_count == 2:
            child_xs = [3.12, 7.20]
        elif child_count == 3:
            child_xs = [0.66, 5.16, 9.66]
        else:
            child_xs = [0.36, 3.54, 6.72, 9.90][:child_count]
        for table_id, x in zip(children, child_xs):
            positions_by_id[table_id] = (x, 5.24)
    else:
        grid = [(0.48, 2.86), (5.16, 2.86), (9.84, 2.86), (0.48, 5.24), (5.16, 5.24), (9.84, 5.24)]
        for table_id, position in zip(ids, grid):
            positions_by_id[table_id] = position

    table_map = {_safe_text(table.get("id"), "", 100): table for table in tables}
    for table_id in ids:
        table = table_map.get(table_id)
        position = positions_by_id.get(table_id)
        if not table or not position:
            continue
        x, y = position
        _erd_table_card(slide, table, x, y, table_w, table_h, accent=C["purple"] if str(database.get("kind")) == "vector" else C["blue"])

    # v5.363 relationship routing: referenced parent -> FK child.
    # Every visible relation gets its own orthogonal lane and its own source/target
    # anchor.  This avoids the old "one shared horizontal bus" where multiple FK
    # lines visually collapsed into a single ambiguous line.
    def line_rect(x: float, y: float, w: float, h: float, color: str):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(max(w, 0.012)), Inches(max(h, 0.012)),
        )
        shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(color); shp.line.fill.background()
        return shp

    def arrow_head(x: float, y: float, w: float, h: float, direction: str, color: str):
        shape_type = {
            "down": MSO_SHAPE.DOWN_ARROW,
            "up": MSO_SHAPE.UP_ARROW,
            "right": MSO_SHAPE.RIGHT_ARROW,
            "left": MSO_SHAPE.LEFT_ARROW,
        }.get(direction, MSO_SHAPE.DOWN_ARROW)
        shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(color); shp.line.fill.background()
        return shp

    visible_relations: list[dict[str, Any]] = []
    for relation in relationships:
        if not isinstance(relation, dict):
            continue
        parent_id = _safe_text(relation.get("to_table"), "", 100)
        child_id = _safe_text(relation.get("from_table"), "", 100)
        if parent_id in positions_by_id and child_id in positions_by_id and parent_id != child_id:
            visible_relations.append({"relation": relation, "parent_id": parent_id, "child_id": child_id})

    visible_relations.sort(key=lambda item: (
        positions_by_id[item["parent_id"]][1],
        positions_by_id[item["parent_id"]][0],
        positions_by_id[item["child_id"]][1],
        positions_by_id[item["child_id"]][0],
        _safe_text((item["relation"] or {}).get("from_column"), "", 80),
    ))

    outgoing_total: dict[str, int] = {}
    incoming_total: dict[str, int] = {}
    for item in visible_relations:
        outgoing_total[item["parent_id"]] = outgoing_total.get(item["parent_id"], 0) + 1
        incoming_total[item["child_id"]] = incoming_total.get(item["child_id"], 0) + 1
    outgoing_seen: dict[str, int] = {}
    incoming_seen: dict[str, int] = {}

    relation_colors = [C["navy2"], C["purple"], C["cyan"], C["green"], C["orange"]]
    cross_relations = [
        item for item in visible_relations
        if positions_by_id[item["parent_id"]][1] < positions_by_id[item["child_id"]][1]
    ]
    # Dedicated central corridor between root and child cards.
    corridor_top = 4.57
    corridor_bottom = 5.13
    lane_count = max(1, len(cross_relations))
    cross_lane_by_key: dict[tuple[str, str, str], float] = {}
    for idx, item in enumerate(cross_relations):
        rel = item["relation"] or {}
        key = (item["parent_id"], item["child_id"], _safe_text(rel.get("from_column"), "", 80))
        cross_lane_by_key[key] = corridor_top + ((idx + 1) / (lane_count + 1)) * (corridor_bottom - corridor_top)

    same_row_seen = 0
    for idx, item in enumerate(visible_relations):
        relation = item["relation"] or {}
        parent_id = item["parent_id"]
        child_id = item["child_id"]
        px, py = positions_by_id[parent_id]
        cx, cy = positions_by_id[child_id]
        color = relation_colors[idx % len(relation_colors)]

        out_index = outgoing_seen.get(parent_id, 0)
        in_index = incoming_seen.get(child_id, 0)
        outgoing_seen[parent_id] = out_index + 1
        incoming_seen[child_id] = in_index + 1
        out_total = max(1, outgoing_total.get(parent_id, 1))
        in_total = max(1, incoming_total.get(child_id, 1))

        # Spread ports across the central 68% of each card so parallel relations
        # do not share the exact same vertical stem.
        parent_anchor = px + table_w * (0.16 + 0.68 * ((out_index + 1) / (out_total + 1)))
        child_anchor = cx + table_w * (0.16 + 0.68 * ((in_index + 1) / (in_total + 1)))

        if py < cy:
            key = (parent_id, child_id, _safe_text(relation.get("from_column"), "", 80))
            lane_y = cross_lane_by_key.get(key, (py + table_h + cy) / 2)
            stem_top = py + table_h
            line_rect(parent_anchor - 0.010, stem_top, 0.020, max(0.02, lane_y - stem_top), color)
            left = min(parent_anchor, child_anchor)
            right = max(parent_anchor, child_anchor)
            line_rect(left, lane_y - 0.010, max(0.02, right - left), 0.020, color)
            arrow_h = max(0.10, cy - lane_y - 0.015)
            arrow_head(child_anchor - 0.055, lane_y, 0.11, arrow_h, "down", color)
        elif py == cy:
            # Same-row FK: route above the row with its own lane instead of drawing
            # through table bodies.  The relation always terminates with an arrow
            # into the FK child card, preserving direction.
            same_row_seen += 1
            lane_y = max(2.68, py - 0.07 - same_row_seen * 0.035)
            line_rect(parent_anchor - 0.010, lane_y, 0.020, max(0.02, py - lane_y), color)
            left = min(parent_anchor, child_anchor)
            right = max(parent_anchor, child_anchor)
            line_rect(left, lane_y - 0.010, max(0.02, right - left), 0.020, color)
            arrow_head(child_anchor - 0.055, lane_y, 0.11, max(0.10, py - lane_y + 0.01), "down", color)
        else:
            # Rare reverse-row relation. Route upward with a unique lane so the
            # diagram remains truthful even when the scoring placed the parent low.
            lane_y = min(py - 0.10, cy + table_h + 0.10 + idx * 0.025)
            line_rect(parent_anchor - 0.010, lane_y, 0.020, max(0.02, py - lane_y), color)
            left = min(parent_anchor, child_anchor)
            right = max(parent_anchor, child_anchor)
            line_rect(left, lane_y - 0.010, max(0.02, right - left), 0.020, color)
            arrow_head(child_anchor - 0.055, cy + table_h, 0.11, max(0.10, lane_y - (cy + table_h)), "up", color)

    message = _safe_text(database.get("message"), "", 180)
    if message and not tables:
        _add_box(slide, 1.28, 3.32, 10.76, 1.64, fill=C["orange_soft"], line=C["orange"])
        _add_text(slide, "Schema 대기", 1.62, 3.62, 2.0, 0.30, size=12, bold=True, color=C["orange"])
        _add_text(slide, message, 3.42, 3.52, 7.92, 0.60, size=9.4, color=C["ink"])
    _add_footer(slide, _safe_text(payload.get("project_name"), "THEANOVA AgentStudio", 100), version, page)
    return slide


def _add_redis_erd_pages(prs: Presentation, payload: dict[str, Any], database: dict[str, Any], page: int, version: str) -> int:
    keys = [item for item in _items(database.get("keys")) if isinstance(item, dict)]
    chunks = [keys[i:i+8] for i in range(0, len(keys), 8)] or [[]]
    for part, chunk in enumerate(chunks, start=1):
        slide = _new_slide(prs)
        suffix = f" · {part}/{len(chunks)}" if len(chunks) > 1 else ""
        _add_title(slide, f"DB ERD · Redis Data Model{suffix}", "Redis는 관계형 ERD 대신 Key Pattern · 역할 · TTL을 논리 데이터 모델로 표현합니다.", "DATABASE · REDIS KEY MODEL")
        _add_large_visual(slide, "cache", 11.54, 0.23, 0.86, 0.88)
        _metric(slide, 0.72, 1.65, 2.62, "Key Patterns", str(len(keys)), "감지/설계된 Redis Key", "purple")
        _metric(slide, 3.54, 1.65, 2.62, "Source", _safe_text(database.get("source"), "INFERENCE", 30), "모델 생성 근거", "cyan")
        _add_box(slide, 6.42, 1.65, 5.46, 1.0, fill=C["panel"], line=C["line"])
        _add_text(slide, "Policy", 6.62, 1.78, 0.8, 0.22, size=8.2, bold=True, color=C["muted"])
        _add_text(slide, _safe_text(database.get("policy"), "Redis runtime/cache model", 180), 7.42, 1.73, 4.18, 0.48, size=8.0, color=C["ink"])
        positions = [(0.72, 2.94), (3.82, 2.94), (6.92, 2.94), (10.02, 2.94), (0.72, 5.00), (3.82, 5.00), (6.92, 5.00), (10.02, 5.00)]
        for item, (x, y) in zip(chunk, positions):
            _add_box(slide, x, y, 2.58, 1.64, fill=C["purple_soft"], line=C["purple"])
            _add_large_visual(slide, "cache", x + 0.10, y + 0.12, 0.72, 0.64)
            _add_text(slide, _safe_text(item.get("key"), "redis:key", 80), x + 0.82, y + 0.10, 1.60, 0.40, size=8.0, bold=True, color=C["ink"])
            _add_text(slide, _safe_text(item.get("purpose"), "Redis Key", 100), x + 0.16, y + 0.76, 2.24, 0.38, size=7.2, color=C["ink"])
            ttl = _safe_text(item.get("ttl"), "TTL 미지정", 42)
            _add_text(slide, f"TTL  {ttl}", x + 0.16, y + 1.18, 2.24, 0.24, size=6.6, color=C["muted"])
        if not chunk:
            _add_text(slide, "Redis 사용은 감지했지만 정적 Key Pattern을 아직 찾지 못했습니다.", 1.55, 3.80, 10.2, 0.54, size=12, bold=True, color=C["muted"], align=PP_ALIGN.CENTER)
        _add_footer(slide, _safe_text(payload.get("project_name"), "THEANOVA AgentStudio", 100), version, page)
        page += 1
    return page


def _add_document_db_pages(prs: Presentation, payload: dict[str, Any], database: dict[str, Any], page: int, version: str) -> int:
    collections = [item for item in _items(database.get("collections")) if isinstance(item, dict)]
    chunks = [collections[i:i+8] for i in range(0, len(collections), 8)] or [[]]
    for part, chunk in enumerate(chunks, start=1):
        slide = _new_slide(prs)
        suffix = f" · {part}/{len(chunks)}" if len(chunks) > 1 else ""
        _add_title(slide, f"DB ERD · {_safe_text(database.get('label'), 'Document DB', 80)}{suffix}", "Document DB는 Collection / Document 구조를 논리 모델로 표현합니다.", "DATABASE · DOCUMENT MODEL")
        _add_large_visual(slide, "cloud", 11.54, 0.23, 0.86, 0.88)
        positions = [(0.72, 1.92), (3.82, 1.92), (6.92, 1.92), (10.02, 1.92), (0.72, 4.38), (3.82, 4.38), (6.92, 4.38), (10.02, 4.38)]
        for item, (x, y) in zip(chunk, positions):
            _add_visual_card(slide, x, y, 2.58, 1.95, title=_safe_text(item.get("name"), "collection", 70), subtitle=_safe_text(item.get("purpose"), "Collection", 90), icon_key="storage", fill=C["cyan_soft"], line=C["cyan"], icon_ratio=0.54, title_size=8.8, subtitle_size=6.7)
        if not chunk:
            _add_text(slide, "Document DB 사용은 감지했지만 Collection 이름을 아직 찾지 못했습니다.", 1.55, 3.58, 10.2, 0.54, size=12, bold=True, color=C["muted"], align=PP_ALIGN.CENTER)
        _add_footer(slide, _safe_text(payload.get("project_name"), "THEANOVA AgentStudio", 100), version, page)
        page += 1
    return page


def _add_db_erd_slides(prs: Presentation, payload: dict[str, Any], page: int, version: str) -> int:
    erd = payload.get("db_erd") or {}
    databases = [item for item in _items(erd.get("databases")) if isinstance(item, dict)]
    if not databases:
        slide = _new_slide(prs)
        _add_title(slide, "DB ERD", "현재 Agent/프로젝트에서 사용하거나 생성되는 DB가 감지되지 않았습니다.", "DATABASE · ERD")
        _add_large_visual(slide, "database", 5.42, 2.32, 2.52, 2.12)
        _add_text(slide, "DB 설계/DDL/Redis Key/Vector Schema가 생성되면 DB별 ERD가 자동으로 추가됩니다.", 2.20, 4.92, 8.92, 0.56, size=11, bold=True, color=C["muted"], align=PP_ALIGN.CENTER)
        _add_footer(slide, _safe_text(payload.get("project_name"), "THEANOVA AgentStudio", 100), version, page)
        return page + 1

    for database in databases:
        kind = str(database.get("kind") or "relational").casefold()
        if kind in {"relational", "vector"}:
            diagram = database.get("diagram") or {}
            tables = [item for item in _items(diagram.get("tables")) if isinstance(item, dict)]
            relationships = [item for item in _items(diagram.get("relationships")) if isinstance(item, dict)]
            # v5.363: dense schemas use fewer tables per slide so PK/FK routes get
            # enough horizontal/vertical space and relation lanes remain traceable.
            relation_count = len([r for r in relationships if isinstance(r, dict)])
            chunk_size = 4 if relation_count >= 36 else (5 if relation_count >= 12 else 6)
            chunks = [tables[i:i+chunk_size] for i in range(0, len(tables), chunk_size)] or [[]]
            for idx, chunk in enumerate(chunks, start=1):
                _add_relational_erd_page(prs, payload, database, chunk, relationships, page, version, idx, len(chunks))
                page += 1
        elif kind == "key-value":
            page = _add_redis_erd_pages(prs, payload, database, page, version)
        elif kind == "document":
            page = _add_document_db_pages(prs, payload, database, page, version)
    return page



def _add_ui_layout_slide(prs, payload: dict[str, Any], page: int, version: str):
    layout = payload.get("ui_layout") or {}
    if not isinstance(layout, dict) or not layout.get("template_id"):
        return False
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "UI / UX Layout", "사용자가 신규 Agent 설계에서 직접 선택한 화면 구조", "AGENT · UI LAYOUT · SELECTED TEMPLATE")
    template_name = _safe_text(layout.get("name") or layout.get("template_name") or layout.get("template_id"), "Selected Layout", 80)
    _add_text(slide, template_name, 0.55, 1.35, 4.6, 0.42, size=20, bold=True, color=C["ink"])
    summary = " · ".join(x for x in [
        _safe_text(layout.get("app_type"), "", 30),
        _safe_text(layout.get("main_layout"), "", 30),
        _safe_text(layout.get("theme_name") if layout.get("theme") == "custom" else layout.get("theme"), "", 40),
        "Responsive" if layout.get("responsive", True) else "Fixed",
    ] if x)
    _add_text(slide, summary, 0.55, 1.80, 5.4, 0.3, size=9, color=C["muted"])

    # Native PowerPoint wireframe so the user can edit every block.
    x, y, w, h = 0.62, 2.30, 7.15, 4.45
    frame = _add_box(slide, x, y, w, h, fill="F7FBFE", line="96BBD2", radius=True)
    header = bool(layout.get("header", True)); sidebar = bool(layout.get("sidebar", False)); footer = bool(layout.get("footer", False)); user_menu = bool(layout.get("user_menu", True))
    top = y + 0.08
    if header:
        _add_box(slide, x+0.10, top, w-0.20, 0.52, fill="DDF1FA", line="B8D6E7", radius=True)
        _add_text(slide, "LOGO", x+0.22, top+0.13, 0.65, 0.22, size=7.5, bold=True, color=C["blue"])
        _add_text(slide, "HOME    SEARCH    DASHBOARD", x+2.00, top+0.13, 2.7, 0.22, size=6.5, color=C["muted"], align=PP_ALIGN.CENTER)
        if user_menu: _add_text(slide, "USER", x+w-0.92, top+0.13, 0.55, 0.22, size=6.5, bold=True, color=C["purple"], align=PP_ALIGN.CENTER)
        top += 0.62
    bottom = y+h-0.10-(0.38 if footer else 0)
    content_x=x+0.10
    if sidebar:
        _add_box(slide, content_x, top, 1.20, bottom-top, fill="EAF5FA", line="C0D8E6", radius=True)
        for i,label in enumerate(["MENU","SEARCH","ORDERS","REPORT","SETTINGS"]):
            _add_text(slide,label,content_x+0.15,top+0.28+i*0.52,0.85,0.22,size=6.4,color=C["navy"])
        content_x += 1.32
    content_w = x+w-0.10-content_x
    _add_box(slide, content_x, top, content_w, 0.38, fill="EAF2FF", line="C7D8F0", radius=True)
    _add_text(slide,"Search / Toolbar",content_x+0.15,top+0.09,content_w-0.30,0.18,size=6.5,color=C["muted"])
    body_y=top+0.52
    cols = 3 if str(layout.get("main_layout")) in {"dashboard","grid"} else (2 if str(layout.get("main_layout")) in {"two_column","three_column"} else 1)
    gap=0.12
    card_w=(content_w-gap*(cols-1))/cols
    components=list(layout.get("components") or [])[:6] or ["Content","Result","Detail"]
    card_count=min(6,max(cols, min(len(components), cols*2)))
    for i in range(card_count):
        col=i%cols; row=i//cols
        card_x=content_x+col*(card_w+gap); card_y=body_y+row*1.22
        _add_box(slide,card_x,card_y,card_w,1.08,fill="FFFFFF",line="C9DCEA",radius=True)
        label=_safe_text(components[i] if i<len(components) else f"Panel {i+1}","Panel",28).replace("_"," ").title()
        _add_text(slide,label,card_x+0.08,card_y+0.12,card_w-0.16,0.22,size=7.2,bold=True,color=C["ink"])
        _add_text(slide,"editable content",card_x+0.08,card_y+0.48,card_w-0.16,0.20,size=5.8,color=C["muted"])
    if footer:
        fy=y+h-0.42
        _add_box(slide,x+0.10,fy,w-0.20,0.30,fill="E8F3F8",line="C5DCE7",radius=True)
        _add_text(slide,"Footer",x+0.22,fy+0.07,w-0.44,0.16,size=6,color=C["muted"],align=PP_ALIGN.CENTER)

    panel_x=8.12
    _add_text(slide,"선택 구성",panel_x,1.45,4.55,0.28,size=12,bold=True,color=C["navy"])
    items=[
        ("App type", layout.get("app_type")), ("Navigation", layout.get("navigation")), ("Main layout", layout.get("main_layout")),
        ("Header", "사용" if header else "없음"), ("Sidebar", "접기 가능" if sidebar and layout.get("sidebar_collapsible") else ("사용" if sidebar else "없음")),
        ("Footer", "사용" if footer else "없음"), ("User menu", _safe_text(layout.get("user_menu_position"),"사용" if user_menu else "없음",30) if user_menu else "없음"),
        ("Theme", layout.get("theme_name") if layout.get("theme") == "custom" else layout.get("theme")), ("Responsive", "Yes" if layout.get("responsive",True) else "No"),
    ]
    yy=1.88
    for label,value in items:
        _add_box(slide,panel_x,yy,4.50,0.43,fill="F8FAFD",line="DCE6EF",radius=True)
        _add_text(slide,label,panel_x+0.13,yy+0.09,1.25,0.20,size=6.5,bold=True,color=C["muted"])
        _add_text(slide,_safe_text(value,"-",60),panel_x+1.45,yy+0.09,2.86,0.20,size=7.2,bold=True,color=C["ink"])
        yy+=0.50
    _add_footer(slide, _safe_text(payload.get("project_name"), "AgentStudio Project", 100), version, page)
    return True


def build_agentstudio_presentation(payload: dict[str, Any], version: str = "5.392") -> tuple[bytes, str]:
    scope = str(payload.get("scope") or "ALL").strip().upper()
    deck_type = str(payload.get("deck_type") or "AGENT").strip().upper()
    if scope not in {"ALL", "WORKFLOW", "RUN", "REPORT", "ARCHITECTURE", "DB_ERD"}:
        scope = "ALL"
    if deck_type not in {"AGENT", "STUDIO"}:
        deck_type = "AGENT"
    if deck_type == "STUDIO":
        scope = "ALL"

    report = payload.get("report") or {}
    style = payload.get("coding_style_report") or {}
    project_name = _safe_text(payload.get("project_name"), "AgentStudio_Project", 100)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    while prs.slides:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    page = 1

    # v5.363 strict deck isolation:
    # - AGENT: only the generated Agent / loaded project.
    # - STUDIO: only THEANOVA AgentStudio itself. No target-project slides are mixed in.
    if deck_type == "STUDIO":
        studio_payload = dict(payload)
        studio_payload["project_name"] = "THEANOVA AgentStudio"
        studio_payload["scope"] = "STUDIO"
        _add_studio_cover(prs, studio_payload, page, version); page += 1
        _add_factory_workflow_slide(prs, studio_payload, page, version); page += 1
        _add_studio_capabilities_slide(prs, studio_payload, page, version); page += 1
        _add_studio_runtime_slide(prs, studio_payload, page, version); page += 1
        _add_studio_governance_slide(prs, studio_payload, page, version); page += 1
        _add_platform_architecture_slide(prs, studio_payload, page, version); page += 1
        _add_platform_foundation_slide(prs, studio_payload, page, version); page += 1
        if _items((studio_payload.get("db_erd") or {}).get("databases")):
            page = _add_db_erd_slides(prs, studio_payload, page, version)

        stream = BytesIO()
        prs.save(stream)
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return stream.getvalue(), f"THEANOVA_AgentStudio_Studio_PPT_전체_{stamp}.pptx"

    # AGENT deck: project/agent data only. AgentStudio factory/platform content is intentionally excluded.
    if scope == "ALL":
        _add_cover(prs, payload, version)
        page += 1
        if _add_ui_layout_slide(prs, payload, page, version):
            page += 1

    if scope in {"ALL", "WORKFLOW"}:
        _add_target_workflow_slide(prs, payload, report, page, version); page += 1

    if scope in {"ALL", "RUN"}:
        _add_execution_slide(prs, payload, report, page, version); page += 1

    if scope in {"ALL", "REPORT"}:
        _add_report_slide(prs, payload, report, style, page, version); page += 1

    if scope in {"ALL", "ARCHITECTURE"}:
        _add_architecture_slide(prs, payload, report, page, version); page += 1
        as_built = report.get("asBuiltArchitecture") or {}
        if _items(as_built.get("components")) or int((as_built.get("scan") or {}).get("source_file_count") or 0) > 0:
            _add_asbuilt_slide(prs, payload, report, page, version); page += 1
        _add_project_stack_slide(prs, payload, report, page, version); page += 1

    if scope == "DB_ERD" or (scope == "ALL" and _items((payload.get("db_erd") or {}).get("databases"))):
        page = _add_db_erd_slides(prs, payload, page, version)

    if len(prs.slides) == 0:
        _add_cover(prs, payload, version)

    stream = BytesIO()
    prs.save(stream)
    safe_name = re.sub(r"[^A-Za-z0-9가-힣._-]+", "_", project_name).strip("._") or "AgentStudio_Project"
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    if scope == "ALL":
        filename = f"{safe_name}_Agent_PPT_전체_{stamp}.pptx"
    else:
        filename = f"{safe_name}_{_SCOPE_LABEL.get(scope, scope)}_{stamp}.pptx"
    return stream.getvalue(), filename

