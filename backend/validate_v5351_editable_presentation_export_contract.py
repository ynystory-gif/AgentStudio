from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
APP = (ROOT / "frontend" / "src" / "App.jsx").read_text(encoding="utf-8")
CSS = (ROOT / "frontend" / "src" / "styles.css").read_text(encoding="utf-8")
ROUTES = (ROOT / "backend" / "app" / "api" / "routes.py").read_text(encoding="utf-8")
SERVICE = (ROOT / "backend" / "app" / "services" / "presentation_export_service.py").read_text(encoding="utf-8")
MAIN = (ROOT / "backend" / "app" / "main.py").read_text(encoding="utf-8")


def require(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


require("AGENTSTUDIO_FRONTEND_VERSION='5.368'" in APP, "frontend version must be 5.368")
require('version="5.368"' in MAIN or "version='5.368'" in MAIN, "FastAPI version must be 5.368")
require('@router.post("/presentation/export")' in ROUTES, "presentation export endpoint missing")
require("build_agentstudio_presentation" in ROUTES, "presentation export service not wired")
require("python-pptx" in (ROOT / "backend" / "requirements.txt").read_text(encoding="utf-8"), "python-pptx dependency missing")

for scope in ("WORKFLOW", "RUN", "REPORT", "ARCHITECTURE"):
    require(f"exportWorkspacePowerPoint('{scope}','AGENT')" in APP, f"{scope} Agent PPT export button missing")
require("exportWorkspacePowerPoint('ALL','AGENT')" in APP, "Agent full PPT button missing")
require("exportWorkspacePowerPoint('ALL','STUDIO')" in APP, "Studio full PPT button missing")

require("workspace-ppt-export-button" in CSS, "PPT export button style missing")
require("native-powerpoint-shapes" in ROUTES, "editable native-shape export contract missing")
require("MSO_SHAPE.ROUNDED_RECTANGLE" in SERVICE, "PowerPoint native shapes not used")
require("_add_architecture_slide" in SERVICE, "editable architecture slide missing")
require("_add_project_stack_slide" in SERVICE, "project adaptive technology/runtime slide missing")
require("_redact" in SERVICE, "PPT export secret redaction missing")

from app.services.presentation_export_service import build_agentstudio_presentation

payload = {
    "scope": "ALL",
    "project_name": "계약 테스트",
    "generated_at": "2026-08-26T13:30:00",
    "workflow_request": "상담형 Agent를 만든다.",
    "workflow_definition": {},
    "coding_style_report": {"checked_files": 7, "pass": 7, "warning": 0, "fail": 0},
    "report": {
        "status": "COMPLETED",
        "testReturncode": 0,
        "testCommand": "python -m pytest",
        "testResult": {"output": "7 passed"},
        "createdFiles": ["backend/app/main.py"],
        "modifiedFiles": ["frontend/src/App.tsx"],
        "debugIteration": 0,
        "debugHistory": [],
        "targetWorkflow": {
            "name": "상담 Workflow",
            "steps": [
                {"label": "입력 수집"},
                {"label": "의도 분석"},
                {"label": "Validator"},
                {"label": "Tool 실행"},
                {"label": "응답 생성"},
            ],
            "branches": [{"if": "tool_required"}],
        },
        "requirementSpec": {
            "goal": "사용자 요청을 분석해 필요한 Tool을 실행하고 답변한다.",
            "acceptance_criteria": ["정상 응답"],
            "constraints": ["Secret 보호"],
        },
        "capabilityPlan": {"capabilities": ["Intent Router", "Validation", "Tool Routing"]},
        "toolMcpPlan": {"decisions": [{"capability": "검색", "execution_type": "MCP"}]},
        "architecture": {
            "components": [
                {"name": "Input Router", "purpose": "입력 처리"},
                {"name": "Agent Core", "purpose": "추론"},
                {"name": "Planner", "purpose": "계획"},
                {"name": "Tool Executor", "purpose": "도구 실행"},
                {"name": "Response Builder", "purpose": "응답"},
            ],
            "interfaces": [{"name": "FastAPI"}],
            "persistence": [{"name": "PostgreSQL"}, {"name": "Redis"}],
            "state": [{"name": "LangGraph State"}],
            "security": [{"name": "Secret Guard"}],
        },
        "asBuiltArchitecture": {
            "components": [{"name": "FastAPI Backend", "status": "detected"}],
            "interfaces": [{"name": "HTTP API"}],
            "persistence": [{"name": "PostgreSQL"}],
            "scan": {"source_file_count": 20},
        },
        "architectureConformance": {"score": 95, "ok": True, "mismatches": []},
        "databasePlan": {"enabled": True, "tables": [{"name": "sessions"}]},
    },
}

content, filename = build_agentstudio_presentation(payload, "5.368")
require(content[:2] == b"PK", "generated pptx is not a zip/OOXML package")
require(filename.endswith(".pptx"), "generated filename must be pptx")
prs = Presentation(BytesIO(content))
require(len(prs.slides) == 7, f"Agent ALL export must contain 7 target-only adaptive slides, got {len(prs.slides)}")
require(sum(len(slide.shapes) for slide in prs.slides) >= 80, "deck should contain editable PowerPoint shapes")
texts = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
for expected in ("상담 Workflow", "실행 결과", "분석 리포트", "Target Project Architecture", "As-Built Architecture", "Project Technology & Runtime"):
    require(expected in texts, f"missing slide content: {expected}")
require("AgentStudio 제작 Workflow" not in texts, "Studio factory workflow must not leak into Agent deck")

print("PASS v5.368 Editable PowerPoint Export contract")
