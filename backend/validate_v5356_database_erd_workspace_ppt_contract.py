from __future__ import annotations

from io import BytesIO
from pathlib import Path
import sys
import tempfile

ROOT = Path(__file__).resolve().parent
PROJECT_ROOT = ROOT.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from app.services.db_erd_service import build_project_db_erd, build_agentstudio_db_erd
from app.services.presentation_export_service import build_agentstudio_presentation


def require(ok: bool, message: str) -> None:
    if not ok:
        raise AssertionError(message)


def slide_texts(blob: bytes) -> list[str]:
    prs = Presentation(BytesIO(blob))
    values: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text)
        values.append("\n".join(parts))
    return values


# A generated Agent can use multiple stores.  Each store must remain a separate
# DB ERD / logical data model instead of being collapsed into one diagram.
with tempfile.TemporaryDirectory(prefix="agentstudio_v5356_erd_") as temp_dir:
    root = Path(temp_dir)
    (root / "postgres_schema.sql").write_text(
        """
        CREATE EXTENSION IF NOT EXISTS vector;
        CREATE TABLE public.customers (
          id BIGSERIAL PRIMARY KEY,
          name TEXT NOT NULL
        );
        CREATE TABLE public.documents (
          id BIGSERIAL PRIMARY KEY,
          customer_id BIGINT NOT NULL REFERENCES public.customers(id),
          content TEXT NOT NULL,
          embedding VECTOR(1536)
        );
        """,
        encoding="utf-8",
    )
    (root / "legacy_mssql.sql").write_text(
        """
        -- Microsoft SQL Server schema
        CREATE TABLE dbo.AuditLog (
          Id INT IDENTITY(1,1) PRIMARY KEY,
          EventName NVARCHAR(200) NOT NULL
        );
        """,
        encoding="utf-8",
    )

    plan = {
        "enabled": True,
        "engine": "postgresql",
        "schema_name": "public",
        "tables": [
            {
                "name": "customers",
                "columns": [
                    {"name": "id", "type": "BIGSERIAL", "primary_key": True, "nullable": False},
                    {"name": "name", "type": "TEXT", "nullable": False},
                ],
            },
            {
                "name": "documents",
                "columns": [
                    {"name": "id", "type": "BIGSERIAL", "primary_key": True, "nullable": False},
                    {"name": "customer_id", "type": "BIGINT", "references": "customers.id", "nullable": False},
                    {"name": "embedding", "type": "VECTOR(1536)", "nullable": True},
                ],
            },
        ],
        "relationships": [
            {"from": "documents.customer_id", "to": "customers.id"},
        ],
        "redis_plan": {
            "enabled": True,
            "keys": [
                {"key": "session:{session_id}", "purpose": "Agent session", "ttl": "1h", "data_type": "hash"},
                {"key": "search:{query_hash}", "purpose": "RAG cache", "ttl": "5m", "data_type": "string/json"},
            ],
        },
    }
    profile = {
        "technologies": ["PostgreSQL", "pgvector", "Redis", "Microsoft SQL Server"],
        "tech_stack": ["PostgreSQL", "pgvector", "Redis", "MSSQL"],
    }
    erd = build_project_db_erd(
        str(root),
        database_plan=plan,
        project_profile=profile,
        workflow_request="PostgreSQL + pgvector + Redis + SQL Server를 사용하는 RAG Agent",
    )

engines = {str(item.get("engine")) for item in erd.get("databases") or []}
require("postgresql" in engines, "PostgreSQL ERD missing")
require("mssql" in engines, "SQL Server ERD must be separated from PostgreSQL")
require("pgvector" in engines, "pgvector logical ERD missing")
require("redis" in engines, "Redis logical data model missing")
require(int((erd.get("summary") or {}).get("database_count") or 0) >= 4, "multiple stores were collapsed")
postgres = next(item for item in erd["databases"] if item.get("engine") == "postgresql")
require(int(postgres.get("table_count") or 0) >= 2, "PostgreSQL tables not detected")
require(int(postgres.get("relationship_count") or 0) >= 1, "PostgreSQL relationship not detected")
redis = next(item for item in erd["databases"] if item.get("engine") == "redis")
require(int(redis.get("key_count") or 0) >= 2, "Redis Key model missing")

# AgentStudio deck has its own DB ERD and is isolated from the current project.
studio_erd = build_agentstudio_db_erd(str(PROJECT_ROOT))
studio_engines = {str(item.get("engine")) for item in studio_erd.get("databases") or []}
require("postgresql" in studio_engines, "AgentStudio PostgreSQL ERD missing")
require(int((studio_erd.get("summary") or {}).get("table_count") or 0) > 0, "AgentStudio schema tables missing")

report = {
    "status": "PROJECT_LOADED",
    "targetWorkflow": {
        "name": "DB ERD Demo Workflow",
        "steps": [
            {"label": "Request", "description": "input"},
            {"label": "Persist", "description": "database"},
            {"label": "Response", "description": "output"},
        ],
        "branches": [],
    },
    "requirementSpec": {"goal": "DB ERD demo"},
    "capabilityPlan": {"capabilities": []},
    "toolMcpPlan": {"decisions": []},
    "architecture": {"components": [], "interfaces": [], "persistence": [], "security": [], "state": [], "infrastructure": []},
    "projectProfile": {"project_type": "RAG_AGENT", "project_type_label": "RAG Agent", "tech_stack": ["PostgreSQL", "pgvector", "Redis", "MSSQL"]},
    "databasePlan": plan,
}
base = {
    "project_name": "db_erd_demo",
    "generated_at": "2026-08-26T17:00:00",
    "workflow_request": "DB ERD demo",
    "workflow_definition": {},
    "coding_style_report": {},
    "report": report,
    "db_erd": erd,
}

erd_blob, erd_name = build_agentstudio_presentation(dict(base, scope="DB_ERD", deck_type="AGENT"), "5.369")
erd_text = "\n".join(slide_texts(erd_blob))
require(erd_blob[:2] == b"PK", "DB ERD page PPTX signature missing")
require("DB ERD" in erd_text, "DB ERD page PPT has no ERD title")
require("PostgreSQL" in erd_text, "DB ERD page PPT missing PostgreSQL")
require("Redis" in erd_text, "DB ERD page PPT missing Redis")
require("_DB_ERD_" in erd_name, "DB ERD page filename missing scope")

agent_blob, _ = build_agentstudio_presentation(dict(base, scope="ALL", deck_type="AGENT"), "5.369")
agent_text = "\n".join(slide_texts(agent_blob))
require("DB ERD" in agent_text and "PostgreSQL" in agent_text, "Agent PPT must include project DB ERD")
require("AgentStudio PostgreSQL / Supabase" not in agent_text, "Studio DB ERD leaked into Agent PPT")

studio_payload = dict(base, scope="ALL", deck_type="STUDIO", db_erd=studio_erd)
studio_blob, _ = build_agentstudio_presentation(studio_payload, "5.369")
studio_text = "\n".join(slide_texts(studio_blob))
require("AgentStudio PostgreSQL / Supabase" in studio_text, "Studio PPT must include AgentStudio DB ERD")
require("db_erd_demo" not in studio_text, "current project leaked into Studio PPT")

app = (PROJECT_ROOT / "frontend" / "src" / "App.jsx").read_text(encoding="utf-8")
panel = (PROJECT_ROOT / "frontend" / "src" / "components" / "database" / "DatabaseErdPanel.tsx").read_text(encoding="utf-8")
routes = (ROOT / "app" / "api" / "routes.py").read_text(encoding="utf-8")
service = (ROOT / "app" / "services" / "db_erd_service.py").read_text(encoding="utf-8")

require("['DB_ERD','DB ERD']" in app, "DB ERD workspace tab missing")
require("exportWorkspacePowerPoint('DB_ERD','AGENT')" in app, "DB ERD page PPT action missing")
require("DatabaseErdPanel" in app and "DatabaseErdPanel" in panel, "DB ERD viewer missing")
require('db_erd:dbErdReport||{}' in app, "Agent/Studio PPT payload must carry ERD snapshot")
require('@router.post("/db-erd/analyze")' in routes, "DB ERD analyze endpoint missing")
require('"DB_ERD"' in routes, "DB ERD presentation scope missing")
require("build_agentstudio_db_erd" in routes, "Studio ERD generation missing")
require("build_project_db_erd" in routes, "Agent/project ERD generation missing")
require("_extract_redis_keys" in service and "_vector_diagram" in service, "Redis/pgvector model inference missing")
require("AGENTSTUDIO_FRONTEND_VERSION='5.369'" in app, "frontend version must be 5.369")

print("PASS v5.369 Database ERD Workspace + PPT contract")
