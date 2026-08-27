from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent
PROJECT_ROOT = ROOT.parent

from app.services.presentation_export_service import build_agentstudio_presentation


def require(condition: bool, message: str) -> None:
    if not condition:
        raise SystemExit("v5.371 contract failed: " + message)


# Version contract.
app = (PROJECT_ROOT / "frontend" / "src" / "App.jsx").read_text(encoding="utf-8")
main = (ROOT / "app" / "main.py").read_text(encoding="utf-8")
routes = (ROOT / "app" / "api" / "routes.py").read_text(encoding="utf-8")
ppt_service = (ROOT / "app" / "services" / "presentation_export_service.py").read_text(encoding="utf-8")
workflow_service = (ROOT / "app" / "services" / "agent_workflow.py").read_text(encoding="utf-8")

require("AGENTSTUDIO_FRONTEND_VERSION='5.371'" in app, "frontend version")
require('version="5.371"' in main or "version='5.371'" in main, "backend version")
require('"version": "5.371"' in routes, "health version")

# Source-level ERD rules.
require("def _set_erd_badge_text" in ppt_service, "horizontal badge helper missing")
require("frame.word_wrap = False" in ppt_service, "PK/FK badge wrapping must be disabled")
require("badge_w = 0.46" in ppt_service, "PK/FK badge width contract missing")
require("cross_lane_by_key" in ppt_service, "relation lane allocator missing")
require("relation_colors" in ppt_service, "relation trace colors missing")
require("chunk_size = 4 if relation_count >= 36 else (5 if relation_count >= 12 else 6)" in ppt_service, "dense ERD auto split missing")

# Generated Agent setup guide rules.
for token in (
    "Ensure-EnvSetupGuides",
    "# DATABASE_URL 입력 방법 (PostgreSQL)",
    "postgresql://사용자:비밀번호@호스트:포트/데이터베이스명",
    "postgresql://postgres:YOUR_PASSWORD@127.0.0.1:5432/postgres",
    "DATABASE_URL 입력 가이드:",
):
    require(token in workflow_service, f"generated SYSTEM_ADMIN/.env guide missing: {token}")
require('if env_key == "DATABASE_URL"' in workflow_service, ".env.example DATABASE_URL guide hook missing")
require('rows.append(f"{env_key}={value}")' in workflow_service, ".env.example key writer missing")

# Build a dense-enough ERD deck and verify actual rendered shape geometry.
def cols(*, fk: str | None = None):
    result = [
        {"name": "id", "data_type": "BIGSERIAL", "primary_key": True},
        {"name": "name", "data_type": "VARCHAR(128)"},
    ]
    if fk:
        result.insert(1, {"name": fk, "data_type": "BIGINT", "foreign_key": True})
    return result


tables = [
    {"id": "public.agents", "schema": "public", "name": "agents", "columns": cols()},
    {"id": "public.workflows", "schema": "public", "name": "workflows", "columns": cols()},
    {"id": "public.agent_versions", "schema": "public", "name": "agent_versions", "columns": cols(fk="agent_id")},
    {"id": "public.agent_features", "schema": "public", "name": "agent_features", "columns": cols(fk="agent_id")},
    {"id": "public.agent_settings", "schema": "public", "name": "agent_settings", "columns": cols(fk="agent_id")},
    {"id": "public.workflow_nodes", "schema": "public", "name": "workflow_nodes", "columns": cols(fk="workflow_id")},
]
relationships = [
    {"from_table": "public.agent_versions", "from_column": "agent_id", "to_table": "public.agents", "to_column": "id"},
    {"from_table": "public.agent_features", "from_column": "agent_id", "to_table": "public.agents", "to_column": "id"},
    {"from_table": "public.agent_settings", "from_column": "agent_id", "to_table": "public.agents", "to_column": "id"},
    {"from_table": "public.workflow_nodes", "from_column": "workflow_id", "to_table": "public.workflows", "to_column": "id"},
]

db_erd = {
    "scope": "AGENT",
    "databases": [
        {
            "engine": "postgresql",
            "label": "PostgreSQL",
            "kind": "relational",
            "source": "AGENT_DATABASE_PLAN",
            "diagram": {"tables": tables, "relationships": relationships},
        }
    ],
}
payload = {
    "project_name": "v5362_erd_demo",
    "generated_at": "2026-08-26T19:30:00",
    "scope": "DB_ERD",
    "deck_type": "AGENT",
    "workflow_request": "",
    "workflow_definition": {},
    "coding_style_report": {},
    "report": {},
    "db_erd": db_erd,
}
blob, _ = build_agentstudio_presentation(payload, "5.371")
require(blob[:2] == b"PK", "PPTX signature")
prs = Presentation(BytesIO(blob))
require(len(prs.slides) >= 1, "ERD slide missing")
slide = prs.slides[0]

badges = []
horizontal_segments = []
for shape in slide.shapes:
    text = ""
    if hasattr(shape, "text"):
        text = str(shape.text or "").strip()
    if text in {"PK", "FK"}:
        badges.append(shape)
    width_in = float(shape.width) / 914400.0
    height_in = float(shape.height) / 914400.0
    top_in = float(shape.top) / 914400.0
    if 4.50 <= top_in <= 5.18 and height_in <= 0.035 and width_in >= 0.10:
        horizontal_segments.append(shape)

require(any(str(getattr(b, "text", "")).strip() == "PK" for b in badges), "PK badge missing")
require(any(str(getattr(b, "text", "")).strip() == "FK" for b in badges), "FK badge missing")
require(all((float(b.width) / 914400.0) >= 0.44 for b in badges), "PK/FK badge width too narrow")
require(all(getattr(b.text_frame, "word_wrap", None) is False for b in badges), "PK/FK badge may wrap vertically")

lane_ys = sorted({round(float(s.top) / 914400.0, 3) for s in horizontal_segments})
require(len(lane_ys) >= 3, f"relation routes still share too few lanes: {lane_ys}")

print("PASS v5.371 ERD Key Badge + Relation Routing + DATABASE_URL Guide contract")
