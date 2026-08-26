from __future__ import annotations

from io import BytesIO
from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from app.services.presentation_export_service import build_agentstudio_presentation


def require(ok: bool, message: str) -> None:
    if not ok:
        raise AssertionError(message)


def texts(blob: bytes) -> list[str]:
    prs = Presentation(BytesIO(blob))
    out: list[str] = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        out.append("\n".join(parts))
    return out


report = {
    "status": "PROJECT_LOADED",
    "targetWorkflow": {
        "name": "Demo Target Workflow",
        "source": "PROJECT_SOURCE_INFERENCE",
        "steps": [
            {"label": "Client Request", "description": "request"},
            {"label": "Project Service", "description": "logic"},
            {"label": "Response", "description": "result"},
        ],
        "branches": [],
    },
    "requirementSpec": {"goal": "Demo target project"},
    "capabilityPlan": {"capabilities": [{"name": "Demo Capability"}]},
    "toolMcpPlan": {"decisions": []},
    "architecture": {
        "components": [
            {"name": "Demo Web UI", "type": "frontend", "description": "target-only-ui"},
            {"name": "Demo API", "type": "backend", "description": "target-only-api"},
        ],
        "interfaces": [{"name": "HTTP API"}],
        "persistence": [],
        "security": [],
        "state": [],
        "infrastructure": [],
    },
    "projectProfile": {"project_type": "WEB_API", "project_type_label": "Web / API Application", "tech_stack": ["React", "FastAPI"], "infrastructure": []},
    "analysisReport": {},
    "createdFiles": [],
    "modifiedFiles": [],
    "debugHistory": [],
    "debugIteration": 0,
    "testCommand": "python -m compileall .",
    "testReturncode": 0,
}
base = {
    "project_name": "demo_target",
    "generated_at": "2026-08-26T15:00:00",
    "workflow_definition": {},
    "coding_style_report": {},
    "report": report,
}

# Agent full deck must contain only target project/agent content.
agent_payload = dict(base, scope="ALL", deck_type="AGENT")
agent_blob, agent_name = build_agentstudio_presentation(agent_payload, "5.368")
agent_text = "\n".join(texts(agent_blob))
require(agent_blob[:2] == b"PK", "Agent PPTX signature missing")
require("Demo Target Workflow" in agent_text, "Agent deck missing target workflow")
require("demo_target" in agent_text, "Agent deck missing target project name")
require("AgentStudio 제작 Workflow" not in agent_text, "Agent deck must not contain Studio factory workflow")
require("THEANOVA AgentStudio 플랫폼 아키텍처" not in agent_text, "Agent deck must not contain Studio platform architecture")
require("AgentStudio Foundation & Infrastructure" not in agent_text, "Agent deck must not contain Studio foundation")
require("_Agent_PPT_전체_" in agent_name, "Agent full filename must be explicit")

# Page-level workflow export is target-only and does not contain Studio content or a full-deck cover.
workflow_payload = dict(base, scope="WORKFLOW", deck_type="AGENT")
workflow_blob, workflow_name = build_agentstudio_presentation(workflow_payload, "5.368")
workflow_slides = texts(workflow_blob)
workflow_text = "\n".join(workflow_slides)
require(len(workflow_slides) == 1, "Agent workflow page export should contain only the target workflow slide")
require("Demo Target Workflow" in workflow_text, "target workflow missing from page export")
require("AgentStudio 제작 Workflow" not in workflow_text, "Studio workflow leaked into page export")
require("_워크플로우_" in workflow_name, "workflow filename missing scope label")

# Studio deck must be self-contained and target-project independent.
studio_payload = dict(base, scope="ALL", deck_type="STUDIO")
studio_blob, studio_name = build_agentstudio_presentation(studio_payload, "5.368")
studio_text = "\n".join(texts(studio_blob))
require(studio_blob[:2] == b"PK", "Studio PPTX signature missing")
require("THEANOVA\nAgentStudio" in studio_text or "AgentStudio" in studio_text, "Studio cover missing")
require("AgentStudio 제작 Workflow" in studio_text, "Studio deck missing factory workflow")
require("AgentStudio Workspace & Core Capabilities" in studio_text, "Studio capability slide missing")
require("AgentStudio Execution & Runtime" in studio_text, "Studio runtime slide missing")
require("AgentStudio Analysis & Governance" in studio_text, "Studio governance slide missing")
require("THEANOVA AgentStudio 플랫폼 아키텍처" in studio_text, "Studio platform architecture missing")
require("AgentStudio Foundation & Infrastructure" in studio_text, "Studio foundation missing")
require("Demo Target Workflow" not in studio_text, "target workflow leaked into Studio deck")
require("demo_target" not in studio_text, "target project name leaked into Studio deck")
require("THEANOVA_AgentStudio_Studio_PPT_전체_" in studio_name, "Studio filename must be explicit")

app = (ROOT.parent / "frontend" / "src" / "App.jsx").read_text(encoding="utf-8")
routes = (ROOT / "app" / "api" / "routes.py").read_text(encoding="utf-8")
service = (ROOT / "app" / "services" / "presentation_export_service.py").read_text(encoding="utf-8")
require("▣ Agent PPT" in app and "▣ Studio PPT" in app, "top export buttons must be split")
require("exportWorkspacePowerPoint('ALL','AGENT')" in app, "Agent PPT action missing")
require("exportWorkspacePowerPoint('ALL','STUDIO')" in app, "Studio PPT action missing")
for scope in ("WORKFLOW", "RUN", "REPORT", "ARCHITECTURE", "DB_ERD"):
    require(f"exportWorkspacePowerPoint('{scope}','AGENT')" in app, f"{scope} page export must be Agent-only")
require('deck_type: str = "AGENT"' in routes, "deck_type request field missing")
require('deck_type == "STUDIO" and scope != "ALL"' in routes, "backend must block page-level Studio exports")
require('if deck_type == "STUDIO"' in service, "Studio deck builder branch missing")
require("AGENTSTUDIO_FRONTEND_VERSION='5.368'" in app, "frontend version must be 5.368")

print("PASS v5.368 Separated Agent / Studio PPT Export contract")
