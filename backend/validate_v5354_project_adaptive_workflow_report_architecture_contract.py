from __future__ import annotations

from io import BytesIO
from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from app.services.project_adaptive_report import _detect_profile
from app.services.presentation_export_service import build_agentstudio_presentation


def require(ok: bool, message: str) -> None:
    if not ok:
        raise AssertionError(message)


def slide_text(blob: bytes) -> str:
    prs = Presentation(BytesIO(blob))
    rows = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                rows.append(shape.text)
    return "\n".join(rows)


# RAG sample: detected technologies must create a RAG-specific workflow/architecture.
rag_files = [
    {
        "relative": "backend/app.py",
        "language": "Python",
        "preview": "from fastapi import FastAPI\nfrom langgraph.graph import StateGraph\nfrom openai import OpenAI\nimport psycopg\nfrom pgvector.psycopg import register_vector\nretriever.similarity_search(query)\n",
    },
    {
        "relative": "frontend/src/App.tsx",
        "language": "React/TypeScript",
        "preview": "import React from 'react'; fetch('/api/search')",
    },
]
rag = _detect_profile(rag_files, "rag_demo")
require(rag["project_type"] == "RAG_AGENT", "RAG source must be classified as RAG_AGENT")
require("Retrieval" in [s["label"] for s in rag["workflow"]["steps"]], "RAG workflow must contain Retrieval")
rag_persistence = [x["label"] for x in rag["architecture"]["persistence"]]
require("PostgreSQL" in rag_persistence and "pgvector" in rag_persistence, "RAG persistence must reflect detected PostgreSQL/pgvector")

# General sample: do not invent Agent/DB/Cloud technologies.
general_files = [
    {
        "relative": "main.py",
        "language": "Python",
        "preview": "def add(a, b):\n    return a + b\n\nprint(add(1, 2))\n",
    }
]
general = _detect_profile(general_files, "simple_python")
require(general["project_type"] == "GENERAL", "simple source must remain GENERAL")
require(not general["architecture"]["persistence"], "GENERAL source must not invent persistence")
require(not general["architecture"]["infrastructure"], "GENERAL source must not invent infrastructure")
require(general["workflow"]["source"] == "PROJECT_SOURCE_INFERENCE", "adaptive workflow source marker missing")
require(general["architecture"]["source"] == "PROJECT_SOURCE_INFERENCE", "adaptive architecture source marker missing")

payload = {
    "scope": "ALL",
    "project_name": "simple_python",
    "project_root": "C:/tmp/simple_python",
    "generated_at": "2026-08-26T14:30:00",
    "workflow_request": "",
    "workflow_definition": {},
    "coding_style_report": {"pass": 0, "warning": 0, "fail": 0, "checked_files": 0},
    "report": {
        "status": "PROJECT_LOADED",
        "targetWorkflow": general["workflow"],
        "requirementSpec": general["requirement_spec"],
        "capabilityPlan": general["capability_plan"],
        "toolMcpPlan": general["tool_mcp_plan"],
        "architecture": general["architecture"],
        "projectProfile": general,
        "analysisReport": general["analysis_report"],
        "createdFiles": [],
        "modifiedFiles": [],
        "debugHistory": [],
        "debugIteration": 0,
        "testCommand": general["execution_baseline"]["test_command"],
        "testReturncode": None,
    },
}
blob, filename = build_agentstudio_presentation(payload, "5.368")
text = slide_text(blob)
require(blob[:2] == b"PK", "PPTX zip signature missing")
require("simple_python · General Software Project Workflow" in text, "project-specific workflow must be exported")
require("Project Technology & Runtime" in text, "adaptive project stack slide missing")
require("THEANOVA AgentStudio 플랫폼 아키텍처" not in text, "fixed AgentStudio platform slide must not be exported for target project")
require("AgentStudio Foundation & Infrastructure" not in text, "fixed AgentStudio foundation slide must not be exported")
for forbidden in ("PostgreSQL", "Redis / Cache", "OpenAI / Ollama", "Kubernetes / Cluster"):
    require(forbidden not in text, f"undetected technology leaked into adaptive PPT: {forbidden}")
require("AgentStudio 제작 Workflow" not in text, "project-load adaptive PPT must prioritize target project workflow")
require(filename.endswith(".pptx"), "PPTX filename missing")

app = (ROOT.parent / "frontend" / "src" / "App.jsx").read_text(encoding="utf-8")
routes = (ROOT / "app" / "api" / "routes.py").read_text(encoding="utf-8")
ppt = (ROOT / "app" / "services" / "presentation_export_service.py").read_text(encoding="utf-8")
require("AGENTSTUDIO_FRONTEND_VERSION='5.368'" in app, "frontend version must be 5.368")
require("/project/adaptive-report" in routes, "adaptive project endpoint missing")
require("build_project_adaptive_report" in routes, "adaptive export refresh missing")
require("loadedProjectAnalysis?.adaptive_report?.workflow" in app, "workflow tab adaptive fallback missing")
require("PROJECT_SOURCE_INFERENCE" in app, "architecture adaptive source handling missing")
require("_add_project_stack_slide" in ppt, "project-specific PPT stack slide missing")

print("PASS v5.368 Project Adaptive Workflow / Report / Architecture contract")
